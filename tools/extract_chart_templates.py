#!/usr/bin/env python3
"""Extract the 4 chart-style templates from a styled .pptx into bundled XML.

generate_pptx.py clones these chartSpace blobs at runtime, so the generator no
longer needs the original template .pptx around. Re-run this only when you want
to restyle the charts: open the template in PowerPoint, adjust colors/fonts/
labels on the 2 demo slides, save, then run:

    python tools/extract_chart_templates.py documents/temp.pptx

Writes tools/chart_templates/{bar,col,donut,stacked}.xml.

Roles (identified by chart type, not slide order):
    bar      bar_clustered, horizontal   → "Total" left chart (MA)
    col      col_clustered,  vertical    → breakdown right charts
    donut    doughnut                    → "Total" left chart (SA≤5)
    stacked  col_percentStacked          → combined breakdown right chart
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

OUT_DIR = Path(__file__).resolve().parent / "chart_templates"


def _bardir(cs):
    bc = cs.find(qn("c:chart") + "/" + qn("c:plotArea") + "/" + qn("c:barChart"))
    if bc is None:
        return None
    d = bc.find(qn("c:barDir"))
    return d.get("val") if d is not None else None


def _grouping(cs):
    bc = cs.find(qn("c:chart") + "/" + qn("c:plotArea") + "/" + qn("c:barChart"))
    if bc is None:
        return None
    g = bc.find(qn("c:grouping"))
    return g.get("val") if g is not None else None


def _is_donut(cs):
    return cs.find(qn("c:chart") + "/" + qn("c:plotArea") + "/" + qn("c:doughnutChart")) is not None


def extract(pptx_path: str) -> None:
    prs = Presentation(pptx_path)
    charts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_chart:
                charts.append(shape.chart._chartSpace)

    roles = {
        "bar":     next((cs for cs in charts if _bardir(cs) == "bar"), None),
        "col":     next((cs for cs in charts if _bardir(cs) == "col"
                         and _grouping(cs) == "clustered"), None),
        "donut":   next((cs for cs in charts if _is_donut(cs)), None),
        "stacked": next((cs for cs in charts if _bardir(cs) == "col"
                         and _grouping(cs) == "percentStacked"), None),
    }
    missing = [k for k, v in roles.items() if v is None]
    if missing:
        sys.exit(f"Could not find template charts for roles: {missing}")

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    for role, cs in roles.items():
        path = OUT_DIR / f"{role}.xml"
        path.write_bytes(etree.tostring(cs, xml_declaration=True,
                                        encoding="UTF-8", standalone=True))
        print(f"  wrote {path}")
    print(f"\nExtracted 4 chart templates -> {OUT_DIR}")


if __name__ == "__main__":
    src = sys.argv[1] if len(sys.argv) > 1 else str(
        Path(__file__).resolve().parent.parent / "documents" / "temp.pptx")
    extract(src)
