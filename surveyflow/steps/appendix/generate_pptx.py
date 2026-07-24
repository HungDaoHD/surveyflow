#!/usr/bin/env python3
"""Generate PowerPoint slides from chart_data.json — template-clone approach.

Instead of building charts from scratch (which never matches the template's
styling), this clones pre-extracted chart-style templates and only swaps in new
data + text. Every visual property — colors, fonts, hidden value axis,
gridlines, gap width, legend, data-label format — is inherited verbatim from the
templates, so the output looks identical to temp.pptx.

The 4 chart styles live as bundled XML in surveyflow/chart_templates/ (extracted
from temp.pptx by extract_chart_templates.py), so the generator is self-contained
— it does NOT need temp.pptx at runtime. To restyle, edit temp.pptx in PowerPoint
and re-run extract_chart_templates.py.

Chart-style roles:
    bar      bar_clustered, horizontal   → "Total" left chart (MA)
    col      col_clustered,  vertical    → breakdown right charts / vertical Total
    donut    doughnut                    → "Total" left chart (SA≤5)
    stacked  col_percentStacked          → combined breakdown right chart

Usage (CLI after pip install surveyflow):
    surveyflow-pptx <chart_data.json> <output.pptx> [options]

Options:
    --templates DIR    Chart-template dir (default: bundled surveyflow/chart_templates)
    --table N          Export ONLY this table index, ungrouped (default: every
                        table marked "is_appendix": true, each as its own
                        group — see _select_tables)
    --start-page N     Starting page number (default: 1)
"""

from __future__ import annotations

import argparse
import copy
import json
import re
import sys
import uuid
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import qn
    from lxml import etree
except ImportError:
    sys.exit("Missing dependency: pip install python-pptx")


# ── Color palette (text only — chart colors come from the template) ───────────

C_NAVY   = RGBColor(0x1F, 0x4E, 0x79)  # title, section labels, Q-prefix   (#1F4E79)
C_ORANGE = RGBColor(0xED, 0x7D, 0x31)  # footer bar                        (#ED7D31)
C_QDESC  = RGBColor(0x40, 0x40, 0x40)  # Q-label description                (#404040)
C_QBASE  = RGBColor(0x7F, 0x7F, 0x7F)  # (N=XX) and page number            (#7F7F7F)


# ── Appendix format "theme" — general vs default (company-branded) ───────────
#
# Two output formats share every geometry/layout FUNCTION in this module
# (they're written generically in terms of an (l, t, w, h) box + fractions,
# never hardcoded pixel math) — only the per-format constants differ: which
# font renders text, which chart color palette is applied, and what color
# section labels / Mean-NPS text use. Bundling those into one `_Theme` object
# and threading it through as a single `theme=` kwarg (default GENERAL_THEME)
# means every existing call site in the "general" render path needed zero
# changes — only the new "default"-format call sites pass `theme=DEFAULT_THEME`
# explicitly.
#
# NOTE: this company-branded format was originally called "Dzung_team"; the
# user/API-facing name is now "default" (also the new default `style`, see
# generate()) and every internal identifier below (DEFAULT_*, _default_*,
# default_template.pptx, chart_templates_default/) was renamed to match —
# "general" is kept fully working but is no longer auto-selected.
class _Theme:
    __slots__ = ("font_name", "chart_palette", "donut_palette", "text_color",
                 "caption_color", "section_font_size", "section_color")

    def __init__(self, *, font_name, chart_palette, donut_palette, text_color,
                 caption_color, section_font_size, section_color):
        self.font_name = font_name
        self.chart_palette = chart_palette
        self.donut_palette = donut_palette
        self.text_color = text_color
        self.caption_color = caption_color
        self.section_font_size = section_font_size
        self.section_color = section_color


# ── Slide geometry — exact EMU values from temp.pptx ─────────────────────────
SW = Emu(12192000)   # 13.333"
SH = Emu(6858000)    # 7.500"

# Title
TITLE_L, TITLE_T = Emu(502920),  Emu(256032)
TITLE_W, TITLE_H = Emu(8686800), Emu(640080)

# Q-label (bottom)
Q_LABEL_L, Q_LABEL_T = Emu(502920),  Emu(5961888)
Q_LABEL_W, Q_LABEL_H = Emu(11521440), Emu(512064)

# Page number
PAGE_L, PAGE_T = Emu(11365992), Emu(6473952)
PAGE_W, PAGE_H = Emu(457200),   Emu(274320)

# Orange footer bar
FOOTER_L, FOOTER_T = Emu(0), Emu(6711696)
FOOTER_W, FOOTER_H = Emu(12188952), Emu(146304)

# Section label height ("Total", "By age")
SECTION_H = Emu(292608)         # 0.32"
SECTION_TO_CHART = Emu(228600)  # gap from section label top to chart top

# Donut ring (plot area) / legend split, as a fraction of the donut chart's
# own bounding box height — ring on top, legend in the remaining band below
# (legendPos="b" in the template). Explicit c:manualLayout, not left to
# PowerPoint's auto-layout: earlier iterations that pinned only PART of this
# (e.g. legend width but not height) each caused a different regression
# (ring inconsistent size, legend stuck top-left, legend overlapping the
# ring for many-choice questions) — but fixing BOTH the ring's and the
# legend's rectangles together, to a plain even split, avoids overlap by
# construction (they can never both claim the same space) and gives every
# question the same ring size regardless of choice/legend-item count.
# Fixing the legend's HEIGHT (rather than leaving it fully automatic) is
# also what makes PowerPoint spread entries across as many columns as fit
# that height, instead of growing a single tall column of rows.
DONUT_RING_FRAC = 0.70

# Width of the legend's manualLayout box, as a fraction of the donut chart's
# own bounding box width, centered horizontally (x = (1 - w) / 2). When this
# was 1.0 (full chart width), PowerPoint spread each row's entries out to
# fill that whole width, leaving very wide gaps between short entries (e.g.
# only 4-6 items across 45% of the slide). Narrowing the box pulls entries
# closer together; centered so the row still looks balanced under the ring.
DONUT_LEGEND_W_FRAC = 0.70

# Vertical center of the ring, as a fraction of the chart's bounding box
# height — used only to position the Mean/NPS overlay text in
# _add_donut_center_text. Derived directly from DONUT_RING_FRAC since the
# ring's plot-area rectangle now spans exactly [0, DONUT_RING_FRAC].
DONUT_CENTER_Y_FRAC = DONUT_RING_FRAC / 2

# Fraction of the stacked-chart's own allocated height reserved (at the top)
# for the 2-line ("7.2" / "-5.0") Mean/NPS label band above each column —
# the chart itself is shrunk by this much so the two never overlap (see
# _build_stacked).
STACK_MEAN_BAND_FRAC = 0.11

# ── Layout A: donut + 100%-stacked (slide2 in template) ──────────────────────
# Donut / breakdown split = 45% / 55% of the available content width.
DT_MARGIN_L = Emu(548640)
DT_GAP      = Emu(228600)
DT_CONTENT_W = int(SW) - 2 * int(DT_MARGIN_L)
_DT_AVAIL_W  = DT_CONTENT_W - int(DT_GAP)
DT_LEFT_W  = Emu(_DT_AVAIL_W * 45 // 100)
DT_RIGHT_W = Emu(_DT_AVAIL_W - int(DT_LEFT_W))
DT_RIGHT_L = Emu(int(DT_MARGIN_L) + int(DT_LEFT_W) + int(DT_GAP))

DT_TOTAL_L, DT_TOTAL_T = DT_MARGIN_L,  Emu(1143000)
DT_TOTAL_W              = DT_LEFT_W
DT_LEFT_L,  DT_LEFT_T   = DT_MARGIN_L,  Emu(1417320)
DT_LEFT_H               = Emu(4023360)
DT_RIGHT_LBL_L, DT_RIGHT_LBL_T = DT_RIGHT_L, Emu(1143000)
DT_RIGHT_LBL_W                  = DT_RIGHT_W
DT_RIGHT_T = Emu(1417320)
DT_RIGHT_H = Emu(4114800)

# ── Layout B: bar charts with N breakdown groups (slide1 in template) ─────────
# Total 40% / Sub-group 60% of content width
BR_TOTAL_L, BR_TOTAL_T = Emu(502920),  Emu(1097280)
BR_TOTAL_W              = Emu(4489704)
BR_LEFT_L,  BR_LEFT_T   = Emu(457200),  Emu(1371600)
BR_LEFT_W,  BR_LEFT_H   = Emu(4535424), Emu(4297680)
BR_RIGHT_LBL_L = Emu(5266944)   # x of right section labels
BR_RIGHT_LBL_W = Emu(6528816)   # width of right section labels
BR_RIGHT_L     = Emu(5221224)   # x of right charts
BR_RIGHT_W     = Emu(6574536)   # width of right charts
BR_RIGHT_FIRST_T = Emu(1097280)  # y of first right section label
BR_RIGHT_BOTTOM  = Emu(5715000)  # bottom of last right chart

# Max breakdown groups shown per slide (Layout B right side)
MAX_GROUPS_PER_SLIDE = 2

# Max stacked-chart breakdown COLUMNS shown per slide for donut_stacked
# (e.g. all banner groups combined can add up to 10+ columns — Area(4) +
# Gender(2) + Age(3) + HHI(4) = 13). Split across continuation slides,
# packing whole banner groups (never splitting one group's columns across
# two slides) — the Total donut is re-rendered identically on every
# continuation slide, only the stack's columns differ.
STACK_MAX_COLS_PER_SLIDE = 10

# ── Layout C: MA bar_horizontal — Total + up to 2 breakdowns, 3 equal columns ─
BH_GAP      = Emu(182880)   # gap between columns (0.2")
BH_MARGIN_L = Emu(457200)   # left margin (mirrors right margin)
BH_COL_W    = Emu((int(SW) - 2 * int(BH_MARGIN_L) - 2 * int(BH_GAP)) // 3)
BH_LBL_T    = Emu(1097280)
BH_CHART_T  = Emu(1371600)
BH_CHART_H  = Emu(4297680)

# Minimum Total-based percent for a choice to appear in an MA cross-tab chart —
# only applied when the question has at least MA_CROSSTAB_MIN_ITEMS choices;
# below that, cross-tab charts just hide 0% items like the Total chart does.
MA_CROSSTAB_MIN_PCT   = 0.10
MA_CROSSTAB_MIN_ITEMS = 6
MA_CROSSTAB_NOTE = (
    f"Note: cross-tab charts only show choices ≥{int(MA_CROSSTAB_MIN_PCT * 100)}% "
    f"of Total (applies to questions with {MA_CROSSTAB_MIN_ITEMS}+ choices)."
)


def _bh_col_l(i: int) -> Emu:
    """Left x-coordinate of the i-th (0-based) equal-width column in Layout C."""
    return Emu(int(BH_MARGIN_L) + i * (int(BH_COL_W) + int(BH_GAP)))


# ── "default" format geometry — exact EMU values from appendix_templates/ ────
# default_template.pptx's 2 demo slides (a company-branded "use_dz" slide
# layout). Only 2 chart archetypes are demonstrated in that template —
# donut+stack (its slide 1) and a 3-equal-column horizontal-bar layout (its
# slide 2) — so the general format's THIRD archetype (bar_vertical / Layout B,
# used for large-choice SA-Nominal questions with no scale) collapses onto
# the same 3-column bar layout as MA's bar_horizontal below, rather than
# inventing an undemonstrated vertical-column style for "default". Title/
# footer/page-number are real placeholders inherited from the bundled
# template's slide layout (idx 13/11/10) — see _default_new_slide — not
# manually positioned textboxes like the general format.
DEFAULT_SW = Emu(9144000)   # 10.000"
DEFAULT_SH = Emu(5143500)   # 5.625"

DEFAULT_FONT = "Segoe UI"
DEFAULT_C_BLACK = RGBColor(0x00, 0x00, 0x00)
DEFAULT_C_GRAY  = RGBColor(0x75, 0x75, 0x75)  # footer text / Mean-NPS caption

# Section/chapter tag (top-left orange box, idx 12) — language-dependent,
# NOT derived from any survey data. Default per lang; override for either one
# via generate(..., default_section_label=...) / --default-section-label.
DEFAULT_SECTION_LABEL_VI = "PHỤ LỤC"
DEFAULT_SECTION_LABEL_EN = "APPENDIX"


def _default_section_label(lang: str) -> str:
    return DEFAULT_SECTION_LABEL_EN if (lang or "vi").strip().lower() == "en" else DEFAULT_SECTION_LABEL_VI

# Layout A: donut + 100%-stacked (slide 1 in default_template.pptx)
DEFAULT_DT_TOTAL_L, DEFAULT_DT_TOTAL_T = Emu(411480),  Emu(857250)
DEFAULT_DT_TOTAL_W                  = Emu(3667316)
DEFAULT_DT_LEFT_L,  DEFAULT_DT_LEFT_T   = Emu(411480),  Emu(1281135)
DEFAULT_DT_LEFT_W,  DEFAULT_DT_LEFT_H   = Emu(3667316), Emu(3017520)
DEFAULT_DT_RIGHT_LBL_L, DEFAULT_DT_RIGHT_LBL_T = Emu(4250245), Emu(857250)
DEFAULT_DT_RIGHT_LBL_W                      = Emu(4482275)
DEFAULT_DT_RIGHT_L = Emu(4250245)
DEFAULT_DT_RIGHT_T = Emu(1402461)
DEFAULT_DT_RIGHT_W = Emu(4482275)
DEFAULT_DT_RIGHT_H = Emu(3017520)

# Layout C: 3 equal-width bar_horizontal columns (slide 2 in default_template.pptx)
DEFAULT_BH_GAP      = Emu(137160)
DEFAULT_BH_MARGIN_L = Emu(342900)
DEFAULT_BH_COL_W    = Emu(2727960)
DEFAULT_BH_LBL_T    = Emu(822960)
DEFAULT_BH_CHART_T  = Emu(1028700)
DEFAULT_BH_CHART_H  = Emu(3403450)


def _default_bh_col_l(i: int) -> Emu:
    """Left x-coordinate of the i-th (0-based) equal-width column in "default" Layout C."""
    return Emu(int(DEFAULT_BH_MARGIN_L) + i * (int(DEFAULT_BH_COL_W) + int(DEFAULT_BH_GAP)))


# ── Namespaces ────────────────────────────────────────────────────────────────
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_RID  = f"{{{_R_NS}}}id"


# ── Chart color palettes (extracted from temp.pptx chart XMLs) ───────────────

# Multi-series palette (breakdown col, stacked) — series 1→6
CHART_PALETTE = [
    "4472C4",  # blue
    "C0504D",  # red
    "9BBB59",  # green
    "95B3D7",  # light blue
    "A6A6A6",  # gray
    "D99694",  # pink/salmon
]

# Donut per-slice palette (matches chart7 in temp.pptx)
# 12 distinct colors — avoids repeats for most SA questions (donut/stack now
# used for all SA regardless of choice count, so >5-choice questions are
# common, e.g. occupation, province, shopping-behavior typology).
DONUT_PALETTE = [
    "4472C4",  # blue
    "95B3D7",  # light blue
    "A6A6A6",  # gray
    "D99694",  # pink/salmon
    "C0504D",  # red
    "9BBB59",  # green
    "8064A2",  # purple
    "F79646",  # orange
    "4BACC6",  # teal
    "808000",  # olive
    "17375E",  # dark navy
    "938953",  # brown/tan
]

# "default"-format palette — the first 5 colors are the ones actually observed in
# the bundled default_template.pptx's own charts (donut dPt fills +
# breakdown-bar series fills); the other 15 are tints/shades of those same 5
# (same brand hues at different lightness — not independently sourced),
# added because a >5-choice donut (e.g. an 8-option MaxDiff-version
# question) was repeating colors once the palette wrapped via
# `palette[i % len(palette)]` (real bug — same slice color reused for a
# different, unrelated choice). Order keeps each tier's 5 brand-color-derived
# entries together and earlier tiers first, so any <=N-choice question's
# colors are byte-identical to before a given extension (only questions with
# MORE choices than the previous palette size are affected by each tier).
# Used for BOTH donut and multi-series roles, unlike the general format's
# separately-curated 12/6-color lists. Re-derive (or replace wholesale) with
# `tools/extract_chart_templates_default.py` against an updated template
# if real examples with more distinct brand colors become available.
DEFAULT_PALETTE = [
    "156082", "0F9ED5", "A6A6A6", "843C0C", "4EA72E",  # brand colors
    "7EA7BA", "7BC9E7", "CECECE", "BB9379", "9DCE8C",  # +45% white tints
    "0F4861", "0B769F", "7C7C7C", "632D09", "3A7D22",  # -25% black shades
    "B8CFD9", "B7E1F2", "E4E4E4", "DAC4B6", "C9E4C0",  # +70% white tints
]

GENERAL_THEME = _Theme(
    font_name=None,  # None → helpers keep their own "Arial" default
    chart_palette=CHART_PALETTE, donut_palette=DONUT_PALETTE,
    text_color=C_NAVY, caption_color=C_QBASE,
    section_font_size=13, section_color=C_NAVY,
)
DEFAULT_THEME = _Theme(
    font_name=DEFAULT_FONT,
    chart_palette=DEFAULT_PALETTE, donut_palette=DEFAULT_PALETTE,
    text_color=DEFAULT_C_BLACK, caption_color=DEFAULT_C_GRAY,
    section_font_size=9.75, section_color=DEFAULT_C_BLACK,
)


# ── Font / text helpers ───────────────────────────────────────────────────────

def _set_run_font(run, font_name: str = "Arial") -> None:
    """Set font (latin + ea + cs) on a text run — "Arial" for the general
    format, "Segoe UI" (DEFAULT_FONT) for "default"."""
    run.font.name = font_name
    rPr = run._r.get_or_add_rPr()
    for tag, pf, cs_val in [(qn("a:ea"), "34", "-122"), (qn("a:cs"), "34", "-120")]:
        el = rPr.find(tag)
        if el is None:
            el = etree.SubElement(rPr, tag)
        el.set("typeface", font_name)
        el.set("pitchFamily", pf)
        el.set("charset", cs_val)


def _set_txbody_margins(tf, anchor: str = "ctr") -> None:
    tf.margin_left = 0
    tf.margin_top = 0
    tf.margin_right = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    # python-pptx's add_textbox() defaults to <a:spAutoFit/> (resize the
    # SHAPE to fit its text) — for short strings like "8.99"/"71.0" this
    # shrinks the box well below the width/position we explicitly computed,
    # which visually collapses a row of evenly-spread per-column labels
    # toward one edge instead of staying centered under their own column.
    # MSO_AUTO_SIZE.NONE (<a:noAutofit/>) keeps the exact box we asked for.
    tf.auto_size = MSO_AUTO_SIZE.NONE
    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    if bodyPr is not None:
        bodyPr.set("anchor", anchor)


def _add_text(slide, left, top, width, height, text, *,
              font_size: int = 12, bold: bool = False,
              color: RGBColor | None = None,
              align=PP_ALIGN.LEFT, anchor: str = "ctr",
              font_name: str = "Arial") -> None:
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    _set_txbody_margins(tf, anchor)
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    _set_run_font(run, font_name)


def _add_q_label(slide, left, top, width, height,
                 question: str, label: str, base) -> None:
    """Q-label: bold navy 'Qxx.' + gray description + gray bold (N=xx)."""
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    _set_txbody_margins(tf, anchor="t")
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.LEFT

    def _run(text, bold=False, color=C_QDESC):
        r = para.add_run()
        r.text = text
        r.font.size = Pt(10)
        r.font.bold = bold
        r.font.color.rgb = color
        _set_run_font(r)

    _run(f"{question}.  ", bold=True, color=C_NAVY)
    _run(label)
    _run(f"  (N={base})", bold=True, color=C_QBASE)


def _set_slide_note(slide, text: str) -> None:
    """Write text into the slide's speaker-notes pane ('Click to add notes').

    No-op (rather than raising) if the notes slide has no body placeholder to
    write into — observed on the bundled "default"-format template, whose notesMaster
    (inherited from a Google Slides export — note the "Google Shape;NNN;p4"
    shape names elsewhere in that template) has no "body" placeholder for
    add_notes_slide() to clone, so notes_text_frame is None. Speaker notes
    here are a diagnostic aid (question type / MA cross-tab caveat), not
    part of the visible appendix, so silently skipping them is safe."""
    if not text:
        return
    tf = slide.notes_slide.notes_text_frame
    if tf is None:
        return
    tf.text = text


def _question_type_label(q: dict) -> str:
    """MA / SA-Ordinal / SA-Nominal, for the speaker-note type tag.

    Driven directly by chart_data's `scale_class` (Claude-classified in
    metadata.json, see CLAUDE.md Step 3c) — SA questions with
    `scale_class == "Ordinal"` are SA-Ordinal, everything else (including
    unclassified SA) is SA-Nominal."""
    answer_type = q.get("answer_type", "")
    if answer_type in ("MA", "Matrix_MA"):
        return "MA"
    if q.get("scale_class") == "Ordinal":
        return "SA-Ordinal"
    return "SA-Nominal"


def _add_rect(slide, left, top, width, height, fill: RGBColor) -> None:
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()


def _add_section_label(slide, left, top, width, text: str, *,
                       height=SECTION_H, theme: "_Theme" = GENERAL_THEME) -> None:
    _add_text(slide, left, top, width, height, text,
              font_size=theme.section_font_size, bold=True,
              color=theme.section_color,
              align=PP_ALIGN.CENTER, anchor="ctr",
              font_name=theme.font_name or "Arial")


# ── Label shortening ──────────────────────────────────────────────────────────

def _shorten_label(text: str) -> str:
    """Remove parenthetical clarifications and piping placeholders from labels, e.g.
    'By using ATM card (Internet banking)' → 'By using ATM card'
    'SHOW ANSWER OF Q13A : {QQ13A/792203/Selected}' → 'SHOW ANSWER OF Q13A :'"""
    result = re.sub(r"\s*\([^)]*\)", "", text)
    result = re.sub(r"\{[^{}]*\}", "", result).strip()
    result = re.sub(r"\s{2,}", " ", result)
    return re.sub(r"[\s:\-]+$", "", result)


def _shorten_labels(items: list, label_key: str = "label") -> list:
    """Shorten a list of choice labels via _shorten_label, but if two items
    collapse to the same shortened text — e.g. two age-banded variants of
    "Married, with children (youngest under 12)" / "...(youngest 12 or
    older)" both become "Married, with children" once the parenthetical
    (the only distinguishing part) is stripped — fall back to the original,
    unshortened text for just those colliding items so choices stay
    distinguishable in the chart/legend."""
    originals = [item[label_key] for item in items]
    shortened = [_shorten_label(t) for t in originals]
    counts: dict = {}
    for s in shortened:
        counts[s] = counts.get(s, 0) + 1
    return [
        orig if counts[short] > 1 else short
        for orig, short in zip(originals, shortened)
    ]


# ── Chart-template cloning ────────────────────────────────────────────────────

def _style_cat_axis(chartspace, *, font_pt: int = 7) -> None:
    """Reduce font size and enable word wrap on category axis labels. No rotation."""
    sz_val = str(font_pt * 100)  # DrawingML unit: 1/100 of a point

    for catAx in chartspace.iter(qn("c:catAx")):
        txPr = catAx.find(qn("c:txPr"))
        if txPr is None:
            txPr = etree.SubElement(catAx, qn("c:txPr"))

        bodyPr = txPr.find(qn("a:bodyPr"))
        if bodyPr is None:
            bodyPr = etree.SubElement(txPr, qn("a:bodyPr"))
        bodyPr.attrib.pop("rot", None)   # remove rotation if previously set
        bodyPr.set("wrap", "square")
        bodyPr.set("anchor", "ctr")

        if txPr.find(qn("a:lstStyle")) is None:
            etree.SubElement(txPr, qn("a:lstStyle"))

        p = txPr.find(qn("a:p"))
        if p is None:
            p = etree.SubElement(txPr, qn("a:p"))
        pPr = p.find(qn("a:pPr"))
        if pPr is None:
            pPr = etree.SubElement(p, qn("a:pPr"))
        defRPr = pPr.find(qn("a:defRPr"))
        if defRPr is None:
            defRPr = etree.SubElement(pPr, qn("a:defRPr"))
        defRPr.set("sz", sz_val)


def _style_legend(chartspace, *, font_pt: int = 7) -> None:
    """Force legend text (item/choice names) to a fixed size across every
    chart type, and force word-wrap so long entries stay fully readable
    within PowerPoint's fixed legend area instead of getting clipped."""
    sz_val = str(font_pt * 100)

    chart_el = chartspace.find(qn("c:chart"))
    if chart_el is None:
        return
    legend = chart_el.find(qn("c:legend"))
    if legend is None:
        return

    txPr = legend.find(qn("c:txPr"))
    if txPr is None:
        txPr = etree.SubElement(legend, qn("c:txPr"))

    bodyPr = txPr.find(qn("a:bodyPr"))
    if bodyPr is None:
        bodyPr = etree.SubElement(txPr, qn("a:bodyPr"))
    bodyPr.set("wrap", "square")
    if txPr.find(qn("a:lstStyle")) is None:
        etree.SubElement(txPr, qn("a:lstStyle"))

    p = txPr.find(qn("a:p"))
    if p is None:
        p = etree.SubElement(txPr, qn("a:p"))
    pPr = p.find(qn("a:pPr"))
    if pPr is None:
        pPr = etree.SubElement(p, qn("a:pPr"))
    defRPr = pPr.find(qn("a:defRPr"))
    if defRPr is None:
        defRPr = etree.SubElement(pPr, qn("a:defRPr"))
    defRPr.set("sz", sz_val)


def _force_pct_labels(chartspace) -> None:
    """Force every data-label number format to 0"%" (values injected as 0-100)."""
    for dl in chartspace.iter(qn("c:dLbls")):
        nf = dl.find(qn("c:numFmt"))
        if nf is None:
            nf = etree.Element(qn("c:numFmt"))
            dl.insert(0, nf)
        nf.set("formatCode", '0"%"')
        nf.set("sourceLinked", "0")



# CT_DLbls / CT_DLbl: elements that must come AFTER c:txPr in document order —
# used to find the correct insertion point when txPr is missing (a plain
# append would violate the schema if any of these are already present).
_DLBL_AFTER_TXPR = (
    "c:dLblPos", "c:showLegendKey", "c:showVal", "c:showCatName",
    "c:showSerName", "c:showPercent", "c:showBubbleSize", "c:separator",
    "c:showLeaderLines", "c:leaderLines", "c:extLst",
)


def _insert_txpr(dl_container):
    """Create a <c:txPr> on a <c:dLbls>/<c:dLbl> element at the schema-correct
    position (before dLblPos/showVal/etc., after idx/numFmt/spPr/any nested
    dLbl), rather than blindly appending at the end."""
    insert_at = len(dl_container)
    for i, child in enumerate(dl_container):
        if etree.QName(child).text in {qn(t) for t in _DLBL_AFTER_TXPR}:
            insert_at = i
            break
    txPr = etree.Element(qn("c:txPr"))
    etree.SubElement(txPr, qn("a:bodyPr"))
    etree.SubElement(txPr, qn("a:lstStyle"))
    dl_container.insert(insert_at, txPr)
    return txPr


def _style_data_labels(chartspace, *, font_pt: int = 6) -> None:
    """Force the percent value shown on every data label (donut slice, bar/
    column end, stacked segment — series-level default AND any per-point
    <c:dLbl> override) to a fixed font size, across every chart type."""
    sz_val = str(font_pt * 100)
    for dl_container in list(chartspace.iter(qn("c:dLbls"))) + list(chartspace.iter(qn("c:dLbl"))):
        txPr = dl_container.find(qn("c:txPr"))
        if txPr is None:
            txPr = _insert_txpr(dl_container)
        p = txPr.find(qn("a:p"))
        if p is None:
            p = etree.SubElement(txPr, qn("a:p"))
        pPr = p.find(qn("a:pPr"))
        if pPr is None:
            pPr = etree.SubElement(p, qn("a:pPr"))
        defRPr = pPr.find(qn("a:defRPr"))
        if defRPr is None:
            defRPr = etree.SubElement(pPr, qn("a:defRPr"))
        defRPr.set("sz", sz_val)


def _hide_zero_value_labels(chartspace) -> None:
    """Delete the data label (on-chart % text / slice) for any point whose
    value is 0, while leaving it in the category/legend cache untouched — so
    the item still appears in the legend, just without a visible 0% slice."""
    chart_el = chartspace.find(qn("c:chart"))
    if chart_el is None:
        return
    plot_area = chart_el.find(qn("c:plotArea"))
    if plot_area is None:
        return
    for ser in plot_area.iter(qn("c:ser")):
        val_el = ser.find(qn("c:val"))
        if val_el is None:
            continue
        numRef = val_el.find(qn("c:numRef"))
        if numRef is None:
            continue
        cache = numRef.find(qn("c:numCache"))
        if cache is None:
            continue
        zero_idxs = sorted(
            int(pt.get("idx")) for pt in cache.findall(qn("c:pt"))
            if float((pt.find(qn("c:v")).text or "0")) == 0
        )
        if not zero_idxs:
            continue
        dLbls = ser.find(qn("c:dLbls"))
        if dLbls is None:
            continue

        # A dLbl per idx must be unique (CT_DLbl: idx, then either <delete> or
        # layout/format children) — reuse the template's existing override for
        # that idx instead of appending a second, conflicting one.
        existing_by_idx = {}
        for dLbl in dLbls.findall(qn("c:dLbl")):
            idx_el = dLbl.find(qn("c:idx"))
            if idx_el is not None:
                existing_by_idx[int(idx_el.get("val"))] = dLbl

        insert_at = 0
        for idx in zero_idxs:
            dLbl = existing_by_idx.get(idx)
            if dLbl is not None:
                for child in list(dLbl):
                    if child.tag != qn("c:idx"):
                        dLbl.remove(child)
                etree.SubElement(dLbl, qn("c:delete")).set("val", "1")
            else:
                dLbl = etree.Element(qn("c:dLbl"))
                etree.SubElement(dLbl, qn("c:idx")).set("val", str(idx))
                etree.SubElement(dLbl, qn("c:delete")).set("val", "1")
                dLbls.insert(insert_at, dLbl)
                insert_at += 1


def _remove_legend(chartspace) -> None:
    chart_el = chartspace.find(qn("c:chart"))
    if chart_el is None:
        return
    legend = chart_el.find(qn("c:legend"))
    if legend is not None:
        chart_el.remove(legend)


def _ensure_legend(chartspace, *, pos: str = "b") -> None:
    """Add a <c:legend> if the template doesn't have one (the "bar" template
    has none, since it's normally used single-series for the Total chart —
    but a multi-series breakdown bar chart needs one to tell series apart).
    No-op if a legend already exists."""
    chart_el = chartspace.find(qn("c:chart"))
    if chart_el is None:
        return
    if chart_el.find(qn("c:legend")) is not None:
        return
    legend = etree.Element(qn("c:legend"))
    etree.SubElement(legend, qn("c:legendPos")).set("val", pos)
    etree.SubElement(legend, qn("c:overlay")).set("val", "0")
    txPr = etree.SubElement(legend, qn("c:txPr"))
    etree.SubElement(txPr, qn("a:bodyPr"))
    etree.SubElement(txPr, qn("a:lstStyle"))
    p = etree.SubElement(txPr, qn("a:p"))
    pPr = etree.SubElement(p, qn("a:pPr"))
    defRPr = etree.SubElement(pPr, qn("a:defRPr"))
    defRPr.set("sz", "900")
    etree.SubElement(defRPr, qn("a:latin")).set("typeface", "Arial")
    etree.SubElement(defRPr, qn("a:cs")).set("typeface", "Arial")
    etree.SubElement(p, qn("a:endParaRPr")).set("lang", "en-US")

    # CT_Chart schema order: ... plotArea, legend?, plotVisOnly?, ...
    plot_area = chart_el.find(qn("c:plotArea"))
    insert_at = list(chart_el).index(plot_area) + 1 if plot_area is not None else len(chart_el)
    chart_el.insert(insert_at, legend)


def _match_series_count(chartspace, n: int) -> None:
    """Make the template's plot have exactly n <c:ser> elements before
    replace_data() runs, by cloning/trimming the template's own series.

    Root-cause fix for a PowerPoint quirk (confirmed via real PowerPoint
    rendering): when replace_data() itself has to add/remove series because
    the new data's series count differs from the template's fixed sample
    count, the chart's LEGEND can render in a different order than the
    series were added in — inconsistently, depending on the count. Matching
    the series count ourselves beforehand means replace_data() only ever
    updates existing series content (never adds/removes any), which renders
    with the legend in the expected add order every time."""
    chart_el = chartspace.find(qn("c:chart"))
    if chart_el is None:
        return
    plot_area = chart_el.find(qn("c:plotArea"))
    if plot_area is None:
        return
    parent = next((child for child in plot_area
                    if child.find(qn("c:ser")) is not None), None)
    if parent is None:
        return
    sers = parent.findall(qn("c:ser"))
    if not sers or n <= 0:
        return
    if len(sers) < n:
        template_ser = sers[-1]
        while len(parent.findall(qn("c:ser"))) < n:
            parent.append(copy.deepcopy(template_ser))
    elif len(sers) > n:
        for extra in parent.findall(qn("c:ser"))[n:]:
            parent.remove(extra)
    for i, ser in enumerate(parent.findall(qn("c:ser"))):
        idx_el = ser.find(qn("c:idx"))
        order_el = ser.find(qn("c:order"))
        if idx_el is not None:
            idx_el.set("val", str(i))
        if order_el is not None:
            order_el.set("val", str(i))


def _match_point_dlbl_count(chartspace) -> None:
    """Extend/trim a single-series chart's per-point <c:dLbl> overrides (e.g.
    donut) to match the actual point count, cloning the last override's
    formatting for any extra points.

    Root-cause fix (confirmed via real PowerPoint rendering): the donut
    template only ships per-point dLbl overrides for its original 5 sample
    slices. Points beyond that (e.g. a 7-choice SA question, now that all SA
    questions use the donut layout) have no override and fall back to the
    series-level defaults — which rendered wrong (showed "0%" instead of the
    real percentage) in real PowerPoint. Giving every point its own override,
    cloned from the same template styling, fixes this for any point count."""
    chart_el = chartspace.find(qn("c:chart"))
    if chart_el is None:
        return
    plot_area = chart_el.find(qn("c:plotArea"))
    if plot_area is None:
        return
    all_sers = [s for child in plot_area for s in child.findall(qn("c:ser"))]
    if len(all_sers) != 1:
        return
    ser = all_sers[0]

    val_el = ser.find(qn("c:val"))
    n = 0
    numRef = val_el.find(qn("c:numRef")) if val_el is not None else None
    if numRef is not None:
        cache = numRef.find(qn("c:numCache"))
        if cache is not None:
            n = len(cache.findall(qn("c:pt")))
    if n == 0:
        return

    dLbls = ser.find(qn("c:dLbls"))
    if dLbls is None:
        return
    dLbl_els = dLbls.findall(qn("c:dLbl"))
    if not dLbl_els:
        return

    if len(dLbl_els) < n:
        template_dLbl = dLbl_els[-1]
        insert_pos = list(dLbls).index(template_dLbl) + 1
        for i in range(len(dLbl_els), n):
            new_dLbl = copy.deepcopy(template_dLbl)
            new_dLbl.find(qn("c:idx")).set("val", str(i))
            dLbls.insert(insert_pos, new_dLbl)
            insert_pos += 1
    elif len(dLbl_els) > n:
        for extra in dLbl_els[n:]:
            dLbls.remove(extra)


def _spPr_solid(hex_color: str):
    """<c:spPr> with solid fill and no border line."""
    spPr = etree.Element(qn("c:spPr"))
    solidFill = etree.SubElement(spPr, qn("a:solidFill"))
    etree.SubElement(solidFill, qn("a:srgbClr")).set("val", hex_color)
    etree.SubElement(etree.SubElement(spPr, qn("a:ln")), qn("a:noFill"))
    return spPr


def _set_ser_spPr(ser, hex_color: str) -> None:
    """Replace or insert <c:spPr> on a series element."""
    existing = ser.find(qn("c:spPr"))
    if existing is not None:
        ser.remove(existing)
    order_el = ser.find(qn("c:order"))
    if order_el is not None:
        ser.insert(list(ser).index(order_el) + 1, _spPr_solid(hex_color))
    else:
        ser.append(_spPr_solid(hex_color))


def _apply_palette(chartspace, palette: list, *, per_point: bool = False,
                   series_colors: list | None = None) -> None:
    """Apply palette colors to chart series.

    per_point=False (bar/col):  single series → 1 color on the series itself.
    per_point=True  (donut):    single series → per-slice dPt colors.
    multi-series (any):         one color per series via spPr — uses
                                 `series_colors[i]` if given (so callers can
                                 keep colors consistent with a paired chart),
                                 otherwise falls back to `palette[i % len]`.
    """
    chart_el = chartspace.find(qn("c:chart"))
    if chart_el is None:
        return
    plot_area = chart_el.find(qn("c:plotArea"))
    if plot_area is None:
        return

    all_sers = [s for child in plot_area for s in child.findall(qn("c:ser"))]
    if not all_sers:
        return

    if len(all_sers) > 1:
        for i, ser in enumerate(all_sers):
            color = series_colors[i] if series_colors else palette[i % len(palette)]
            _set_ser_spPr(ser, color)
        return

    ser = all_sers[0]
    if not per_point:
        # Bar/col single series: uniform color
        _set_ser_spPr(ser, series_colors[0] if series_colors else palette[0])
    else:
        # Donut: color each slice individually via dPt
        for dPt in ser.findall(qn("c:dPt")):
            ser.remove(dPt)
        val_el = ser.find(qn("c:val"))
        n = 0
        if val_el is not None:
            numRef = val_el.find(qn("c:numRef"))
            if numRef is not None:
                cache = numRef.find(qn("c:numCache"))
                if cache is not None:
                    n = len(cache.findall(qn("c:pt")))
        if n == 0:
            return
        children = list(ser)
        insert_pos = next(
            (i for i, c in enumerate(children)
             if c.tag in (qn("c:cat"), qn("c:val"))),
            len(children),
        )
        for i in range(n):
            dPt = etree.Element(qn("c:dPt"))
            etree.SubElement(dPt, qn("c:idx")).set("val", str(i))
            dPt.append(_spPr_solid(palette[i % len(palette)]))
            ser.insert(insert_pos + i, dPt)


def _set_donut_layout(chartspace, *, ring_frac: float) -> None:
    """Force explicit c:manualLayout so the ring (plot area) occupies the top
    `ring_frac` of the chart's bounding box and the legend occupies the
    remaining bottom band — see DONUT_RING_FRAC for rationale. Both
    rectangles are set together (never just one), so they can't overlap."""
    chart_el = chartspace.find(qn("c:chart"))
    if chart_el is None:
        return

    plot_area = chart_el.find(qn("c:plotArea"))
    if plot_area is not None:
        old_layout = plot_area.find(qn("c:layout"))
        if old_layout is not None:
            plot_area.remove(old_layout)
        layout = etree.Element(qn("c:layout"))
        manual = etree.SubElement(layout, qn("c:manualLayout"))
        etree.SubElement(manual, qn("c:layoutTarget")).set("val", "inner")
        etree.SubElement(manual, qn("c:xMode")).set("val", "edge")
        etree.SubElement(manual, qn("c:yMode")).set("val", "edge")
        etree.SubElement(manual, qn("c:x")).set("val", "0")
        etree.SubElement(manual, qn("c:y")).set("val", "0")
        etree.SubElement(manual, qn("c:w")).set("val", "1")
        etree.SubElement(manual, qn("c:h")).set("val", str(ring_frac))
        # CT_PlotArea: layout is always the first optional child.
        plot_area.insert(0, layout)

    legend = chart_el.find(qn("c:legend"))
    if legend is not None:
        old_layout = legend.find(qn("c:layout"))
        if old_layout is not None:
            legend.remove(old_layout)
        legend_x = round((1 - DONUT_LEGEND_W_FRAC) / 2, 4)
        layout = etree.Element(qn("c:layout"))
        manual = etree.SubElement(layout, qn("c:manualLayout"))
        etree.SubElement(manual, qn("c:xMode")).set("val", "edge")
        etree.SubElement(manual, qn("c:yMode")).set("val", "edge")
        etree.SubElement(manual, qn("c:x")).set("val", str(legend_x))
        etree.SubElement(manual, qn("c:y")).set("val", str(round(ring_frac, 4)))
        etree.SubElement(manual, qn("c:w")).set("val", str(DONUT_LEGEND_W_FRAC))
        etree.SubElement(manual, qn("c:h")).set("val", str(round(1 - ring_frac, 4)))
        # CT_Legend order: legendPos?, legendEntry*, layout?, overlay?, spPr?, txPr?
        entries = legend.findall(qn("c:legendEntry"))
        anchor = entries[-1] if entries else legend.find(qn("c:legendPos"))
        insert_at = list(legend).index(anchor) + 1 if anchor is not None else 0
        legend.insert(insert_at, layout)


def _clone_chart(slide, tmpl_chartspace, pptx_type, l, t, w, h, chart_data,
                 *, drop_legend: bool = False, per_point: bool = False,
                 palette: list | None = None,
                 cat_font_pt: int = 7, legend_n: int | None = None,
                 hide_zero_labels: bool = False, series_colors: list | None = None,
                 donut_ring_frac: float | None = None):
    """Add a chart, replace its XML with a deep copy of the template chartSpace,
    then inject `chart_data` via replace_data() so template styling is kept.

    Item text (category axis + legend) and data-label percent text are
    always forced to a fixed size (7pt / 6pt) for every chart type — see
    _style_cat_axis / _style_legend / _style_data_labels. Both are no-ops
    when the chart has no category axis (donut) or no legend, so calling
    them unconditionally is safe."""
    gf = slide.shapes.add_chart(pptx_type, l, t, w, h, chart_data)
    cp = gf.chart_part

    new_ext = cp._element.find(qn("c:externalData"))
    new_rid = new_ext.get(_RID) if new_ext is not None else None

    cs = copy.deepcopy(tmpl_chartspace)
    t_ext = cs.find(qn("c:externalData"))
    if t_ext is not None and new_rid is not None:
        t_ext.set(_RID, new_rid)

    if drop_legend:
        _remove_legend(cs)
    _match_series_count(cs, len(chart_data))

    cp._element = cs
    cp.__dict__.pop("chart", None)
    cp.__dict__.pop("chart_workbook", None)

    cp.chart.replace_data(chart_data)
    _force_pct_labels(cp._element)
    _style_cat_axis(cp._element, font_pt=cat_font_pt)
    if legend_n is not None:
        _ensure_legend(cp._element)
    _style_legend(cp._element)
    if per_point:
        _match_point_dlbl_count(cp._element)
    if hide_zero_labels:
        _hide_zero_value_labels(cp._element)
    _style_data_labels(cp._element)
    _apply_palette(cp._element, palette or CHART_PALETTE, per_point=per_point,
                   series_colors=series_colors)
    if donut_ring_frac is not None:
        _set_donut_layout(cp._element, ring_frac=donut_ring_frac)
    return cp.chart


def _v100(percents: dict, codes: list) -> list:
    """Fractions (0-1) to percentages (0-100) for the given choice codes."""
    return [round(percents.get(c, 0.0) * 100.0, 4) for c in codes]


def _mean_nps_lines(data: dict) -> list[str]:
    """['Mean: 7.2'] / ['Mean: 7.2', 'NPS: -5.0'] for a total/breakdown-column
    dict that carries Step 3c's `mean`/`nps` stat values — Mean and NPS are
    always separate lines (never joined on one line), and this returns []
    when the question has no `mean` (e.g. Nominal or MA), so it's always
    safe to call unconditionally."""
    mean_v = data.get("mean")
    if mean_v is None:
        return []
    lines = [f"Mean: {mean_v:g}"]
    nps_v = data.get("nps")
    if nps_v is not None:
        lines.append(f"NPS: {nps_v:.1f}")
    return lines


def _mean_nps_compact(data: dict) -> list[str]:
    """['9.24'] / ['9.24', '82.8'] — bare numbers, one per line (Mean then
    NPS), used for the stack chart's per-column top labels. A single
    "Mean" / "NPS" caption to the left labels the two rows once (see
    _build_stacked), rather than repeating "Mean:"/"NPS:" prefixes on
    every column."""
    mean_v = data.get("mean")
    if mean_v is None:
        return []
    lines = [f"{mean_v:g}"]
    nps_v = data.get("nps")
    if nps_v is not None:
        lines.append(f"{nps_v:.1f}")
    return lines


def _col_label_with_base(col: dict) -> str:
    """Breakdown column label + its own base, e.g. 'Male (N=120)'.

    Mean/NPS (Step 3c Ordinal questions) is NOT appended here — it's shown
    as a data label sitting on top of that column's stacked bar instead
    (see _set_stack_top_labels), so the x-axis stays a plain category name."""
    return f"{col.get('label', '')} (N={col.get('base', 0)})"


# Meta codes that are never part of an ordinal scale itself, even when they
# appear alongside one — per this project's own FT-coding convention (see
# CLAUDE.md Workflow C: 98 = "Others", 99 = "No answer"). Sorting these
# purely by numeric code value would put them at the very top of a 1-5
# scale (since 98/99 > 5), which is never the intended meaning.
_ORDINAL_META_CODES = {98, 99}


def _code_int(choice: dict) -> int | None:
    """Parse a single choice's code as int, or None if it isn't numeric."""
    try:
        return int(choice["code"])
    except (ValueError, TypeError):
        return None


def _sort_scale_desc(q: dict, choices: list) -> list:
    """SA-Ordinal questions (`scale_class == "Ordinal"`, Claude-classified in
    metadata.json — see CLAUDE.md Step 3c): sort choices highest scale point
    to lowest. SA-Nominal questions are returned unchanged.

    Which code is "highest" isn't always the max code value — e.g. a
    frequency scale coded 1=Almost every day … 7=Less than once a month has
    code 1 at the high-intensity end. `scale_high_code` (also
    Claude-classified, see CLAUDE.md Step 3c) names the code at the high end
    explicitly; sort direction is derived from whether that's the min or max
    code among the real scale choices (excluding meta codes below).
    Defaults to "max code = high end" (descending by code) when
    `scale_high_code` isn't set.

    Choices with a non-numeric code, or a meta code (_ORDINAL_META_CODES),
    are excluded from the scale sort and always appended at the end (in
    their original relative order) — sorting them in with the real scale
    values would place e.g. "Others"/"No answer" at the top purely because
    98/99 is numerically larger than the scale's real range."""
    if q.get("scale_class") != "Ordinal":
        return choices

    scale_choices, meta_choices = [], []
    for c in choices:
        code = _code_int(c)
        if code is None or code in _ORDINAL_META_CODES:
            meta_choices.append(c)
        else:
            scale_choices.append(c)
    if not scale_choices:
        return choices

    codes = [_code_int(c) for c in scale_choices]
    high_code = q.get("scale_high_code")
    reverse = True
    if high_code is not None:
        try:
            reverse = int(high_code) != min(codes)
        except (ValueError, TypeError):
            pass
    scale_choices.sort(key=_code_int, reverse=reverse)
    return scale_choices + meta_choices


def _order_sa_choices(q: dict, choices: list) -> list:
    """Order SA (donut_stacked) choices for both the donut and its paired
    stacked breakdown chart, so the two always agree:

    - NPS questions (`is_nps`, see CLAUDE.md Step 3c): never re-sorted —
      chart_data_exporter already emits exactly 3 choices (Promoters,
      Passives, Detractors) in that fixed order, keyed by synthetic codes
      "0"/"1"/"2". Sorting by scale value here would be meaningless (and
      actively wrong) since those aren't real scale codes.
    - Range/bucket groups (`preserve_order` — NUM `num_quantile` or an
      all-hidden "choices" config, see chart_data_exporter.py): also never
      re-sorted, for the same reason as NPS — the groups already come out
      in a meaningful order (e.g. ascending age bands "20-28"/"29-34"/...)
      that Total-percent sorting would scramble.
    - SA-Ordinal: sort by scale order descending (see _sort_scale_desc).
    - SA-Nominal: sort by Total percent descending, tie-broken by ascending
      code (stable/deterministic tie-break, rather than silently depending
      on whatever order the choices happened to arrive in — which mattered
      in practice for questions with several 0%/tied choices)."""
    if q.get("is_nps") or q.get("preserve_order"):
        return choices
    if q.get("scale_class") == "Ordinal":
        return _sort_scale_desc(q, choices)
    total_pct = q.get("total", {}).get("percents", {})

    def _key(c):
        code = _code_int(c)
        return (-total_pct.get(c["code"], 0), code if code is not None else 0)

    return sorted(choices, key=_key)


# ── Per-chart builders (each clones the matching template chart) ──────────────

def _build_total_bar(slide, q, tmpl, l, t, w, h, *,
                     theme: "_Theme" = GENERAL_THEME) -> None:
    """Horizontal bar, single series, sorted ascending (largest at top). Hides 0% items."""
    choices = q["choices"]
    pcts = q.get("total", {}).get("percents", {})
    choices = [c for c in choices if pcts.get(c["code"], 0) > 0]
    if not choices:
        return
    codes = [c["code"] for c in choices]
    labels = _shorten_labels(choices)
    order = sorted(range(len(codes)), key=lambda i: pcts.get(codes[i], 0))
    labels = [labels[i] for i in order]
    codes  = [codes[i]  for i in order]
    cd = CategoryChartData()
    cd.categories = labels
    cd.add_series("%", _v100(pcts, codes))
    _clone_chart(slide, tmpl["bar"], XL_CHART_TYPE.BAR_CLUSTERED, l, t, w, h, cd,
                 palette=theme.chart_palette)


def _build_total_col(slide, q, tmpl, l, t, w, h) -> None:
    """Vertical column, single series. Hides 0% items. SA-Likert questions
    (scale detected) sort descending by scale point; everything else keeps
    natural order."""
    choices = q["choices"]
    pcts = q.get("total", {}).get("percents", {})
    choices = [c for c in choices if pcts.get(c["code"], 0) > 0]
    if not choices:
        return
    choices = _sort_scale_desc(q, choices)
    codes  = [c["code"] for c in choices]
    labels = _shorten_labels(choices)
    cd = CategoryChartData()
    cd.categories = labels
    cd.add_series("Total", _v100(pcts, codes))
    _clone_chart(slide, tmpl["col"], XL_CHART_TYPE.COLUMN_CLUSTERED, l, t, w, h, cd,
                 drop_legend=True)


def _build_breakdown_col(slide, q, group, tmpl, l, t, w, h) -> None:
    """Vertical clustered columns — one series per breakdown column. Hides 0%
    items. SA-Likert questions (scale detected) sort descending by scale
    point, matching the Total chart; everything else keeps natural order."""
    choices = q["choices"]
    cols = group["columns"]
    choices = [c for c in choices
               if any(col.get("percents", {}).get(c["code"], 0) > 0 for col in cols)]
    if not choices:
        return
    choices = _sort_scale_desc(q, choices)
    codes  = [c["code"] for c in choices]
    labels = _shorten_labels(choices)
    cd = CategoryChartData()
    cd.categories = labels
    for col in cols:
        cd.add_series(col["label"], _v100(col.get("percents", {}), codes))
    _clone_chart(slide, tmpl["col"], XL_CHART_TYPE.COLUMN_CLUSTERED, l, t, w, h, cd)


def _build_breakdown_bar(slide, q, group, tmpl, l, t, w, h, *,
                         theme: "_Theme" = GENERAL_THEME) -> None:
    """Horizontal clustered bars — one series per breakdown column, ordered like
    the Total chart (largest at top). All MA bar charts hide 0% items; for
    questions with >= MA_CROSSTAB_MIN_ITEMS choices, the MA_CROSSTAB_MIN_PCT
    cutoff is applied on top of that."""
    choices = q["choices"]
    cols = group["columns"]
    total_pct = q.get("total", {}).get("percents", {})
    if len(choices) >= MA_CROSSTAB_MIN_ITEMS:
        choices = [c for c in choices if total_pct.get(c["code"], 0) >= MA_CROSSTAB_MIN_PCT]
    else:
        choices = [c for c in choices if total_pct.get(c["code"], 0) > 0]
    if not choices:
        return
    codes  = [c["code"] for c in choices]
    labels = _shorten_labels(choices)
    order = sorted(range(len(codes)), key=lambda i: total_pct.get(codes[i], 0))
    codes  = [codes[i]  for i in order]
    labels = [labels[i] for i in order]
    cd = CategoryChartData()
    cd.categories = labels
    for col in cols:
        cd.add_series(col["label"], _v100(col.get("percents", {}), codes))
    _clone_chart(slide, tmpl["bar"], XL_CHART_TYPE.BAR_CLUSTERED, l, t, w, h, cd,
                 legend_n=len(cols) if len(cols) > 1 else None,
                 palette=theme.chart_palette)


def _build_breakdown_stacked(slide, q, group, tmpl, l, t, w, h, *,
                             horizontal: bool = False) -> None:
    """100%-stacked breakdown chart for questions with many choices (> STACKED_THRESHOLD).
    Categories = subgroup columns (e.g. Male / Female).
    Series     = choices (items), 0% filtered.
    horizontal=True  → BAR_STACKED_100  (for MA bar_horizontal questions)
    horizontal=False → COLUMN_STACKED_100
    """
    choices = q["choices"]
    cols = group["columns"]
    choices = [c for c in choices
               if any(col.get("percents", {}).get(c["code"], 0) > 0 for col in cols)]
    if not choices:
        return
    cd = CategoryChartData()
    cd.categories = [_col_label_with_base(col) for col in cols]
    choice_labels = _shorten_labels(choices)
    for choice, choice_label in zip(choices, choice_labels):
        cd.add_series(
            choice_label,
            [round(col.get("percents", {}).get(choice["code"], 0.0) * 100.0, 4)
             for col in cols],
        )
    if horizontal:
        _clone_chart(slide, tmpl["bar"], XL_CHART_TYPE.BAR_STACKED_100,
                     l, t, w, h, cd)
    else:
        _clone_chart(slide, tmpl["stacked"], XL_CHART_TYPE.COLUMN_STACKED_100,
                     l, t, w, h, cd)


def _add_multiline_text(slide, left, top, width, height, lines: list[str], *,
                        font_size: int = 9, bold: bool = True,
                        color: RGBColor | None = None,
                        font_name: str = "Arial") -> None:
    """Small centered text box with each string in `lines` as its own
    paragraph (never joined with a separator onto one line)."""
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    _set_txbody_margins(tf, anchor="ctr")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = line
        r.font.size = Pt(font_size)
        r.font.bold = bold
        if color:
            r.font.color.rgb = color
        _set_run_font(r, font_name)


def _add_donut_center_text(slide, l, t, w, h, lines: list[str], *,
                           theme: "_Theme" = GENERAL_THEME) -> None:
    """Small Mean/NPS text box sitting inside the donut's hole (holeSize=55%
    of radius, see donut.xml). The ring's plot area is manually pinned to
    the top DONUT_RING_FRAC of the chart box (see _set_donut_layout), so
    its vertical center is exactly DONUT_CENTER_Y_FRAC — no longer a
    render-and-check heuristic."""
    box_w = Emu(int(w) * 2 // 5)
    box_h = Emu(int(h) // 6)
    box_l = Emu(int(l) + (int(w) - int(box_w)) // 2)
    box_t = Emu(int(t) + int(int(h) * DONUT_CENTER_Y_FRAC) - int(box_h) // 2)
    _add_multiline_text(slide, box_l, box_t, box_w, box_h, lines,
                        font_size=9, bold=True, color=theme.text_color,
                        font_name=theme.font_name or "Arial")


def _build_donut(slide, q, tmpl, l, t, w, h, *,
                 theme: "_Theme" = GENERAL_THEME) -> None:
    """Donut chart for SA≤5. Every choice (incl. 0%) stays in the legend, but
    0%-value slices/labels are hidden on the chart itself. Order: scale
    questions descending by scale point, non-scale by Total percent
    descending (see _order_sa_choices)."""
    choices = q["choices"]
    pcts = q.get("total", {}).get("percents", {})
    if not choices:
        return
    choices = _order_sa_choices(q, choices)
    codes  = [c["code"] for c in choices]
    labels = _shorten_labels(choices)
    cd = CategoryChartData()
    cd.categories = labels
    cd.add_series("Total", _v100(pcts, codes))
    _clone_chart(slide, tmpl["donut"], XL_CHART_TYPE.DOUGHNUT, l, t, w, h, cd,
                 per_point=True, palette=theme.donut_palette, legend_n=len(labels),
                 hide_zero_labels=True, donut_ring_frac=DONUT_RING_FRAC)
    center_lines = _mean_nps_lines(q.get("total", {}))
    if center_lines:
        _add_donut_center_text(slide, l, t, w, h, center_lines, theme=theme)


def _add_stack_top_labels(slide, l, t, w, band_h, lines_per_col: list[list[str]], *,
                          theme: "_Theme" = GENERAL_THEME) -> None:
    """A small 2-line ("7.2" / "-5.0") Mean/NPS textbox per column, each one
    centered on that column's own horizontal slot (w / n, matching how
    PowerPoint centers a category's bar within its own equal-width slot —
    same assumption category axes use by default), inside the reserved top
    band (l, t, w, band_h). The caller has already shrunk the actual chart
    to start below this band, so there's no dependency on estimating where
    the real chart's plot area happens to render; the two literally cannot
    overlap.

    A single 2-line "Mean" / "NPS" caption is added once, to the left of
    the first column (in the gap before the chart's own left edge), lined
    up row-for-row with the per-column values — rather than repeating
    "Mean:"/"NPS:" labels on every column."""
    n = len(lines_per_col)
    if n == 0:
        return
    col_w = int(w) // n
    box_w = Emu(int(col_w * 0.92))
    box_h = Emu(int(band_h))

    # Only label "NPS" if at least one column actually has an nps value —
    # a Mean-only Ordinal question (no NPS stat) should show just "Mean".
    has_nps = any(len(lines) > 1 for lines in lines_per_col)
    caption = ["Mean", "NPS"] if has_nps else ["Mean"]
    caption_l = Emu(max(0, int(l) - int(box_w)))
    font_name = theme.font_name or "Arial"
    _add_multiline_text(slide, caption_l, Emu(int(t)), box_w, box_h,
                        caption, font_size=7, bold=True, color=theme.caption_color,
                        font_name=font_name)

    for i, lines in enumerate(lines_per_col):
        if not lines:
            continue
        col_l = int(l) + i * col_w
        box_l = Emu(col_l + (col_w - int(box_w)) // 2)
        _add_multiline_text(slide, box_l, Emu(int(t)), box_w, box_h, lines,
                            font_size=7, bold=True, color=theme.text_color,
                            font_name=font_name)


def _build_stacked(slide, q, breakdowns, tmpl, l, t, w, h, *,
                   theme: "_Theme" = GENERAL_THEME) -> None:
    """100%-stacked columns — series per choice, using the same donut_palette
    colors per choice so the two charts agree. Shows every choice (incl.
    those that are 0% everywhere) so the legend always matches the donut's
    full list — donut's post-sort order (_order_sa_choices) is the source of
    truth for both the stack's item order (within a column) and its legend
    order.

    Series are always added in the same order as the donut's sorted choices
    (SA-Nominal / SA-Ordinal / Matrix-SA alike) — no reversal. PowerPoint's
    legend for THIS chart can render in reverse of series-add order when
    labels are long enough to force a single-column legend layout (observed
    on long-label questions like S10/A5/B6) — accepted as a known edge case
    rather than compensated for, since a label-length-based reversal rule
    was tried and rejected: it fixed those long-label cases but flipped the
    order on other (shorter-label) Nominal questions that were already
    correct without any reversal."""
    choices = q["choices"]
    all_cols, xlabels = [], []
    for bd in breakdowns:
        for col in bd["columns"]:
            all_cols.append(col)
            xlabels.append(_col_label_with_base(col))
    if not all_cols or not choices:
        return
    choices = _order_sa_choices(q, choices)
    palette = theme.donut_palette
    color_map = {c["code"]: palette[i % len(palette)]
                 for i, c in enumerate(choices)}
    shortened = {c["code"]: lbl for c, lbl in zip(choices, _shorten_labels(choices))}
    cd = CategoryChartData()
    cd.categories = xlabels
    series_colors = []
    for choice in choices:
        cd.add_series(
            shortened[choice["code"]],
            [round(col.get("percents", {}).get(choice["code"], 0.0) * 100.0, 4)
             for col in all_cols],
        )
        series_colors.append(color_map[choice["code"]])

    # Mean/NPS (Step 3c Ordinal questions) sits above each column, in a
    # reserved band carved out of the TOP of this function's own (l, t, w, h)
    # allocation — the chart itself is shrunk/shifted down to make room, so
    # the label band and the chart body never overlap.
    mean_lines = [_mean_nps_compact(col) for col in all_cols]
    if any(mean_lines):
        chart_t = Emu(int(t) + int(int(h) * STACK_MEAN_BAND_FRAC))
        chart_h = Emu(int(h) - int(int(h) * STACK_MEAN_BAND_FRAC))
        _add_stack_top_labels(slide, l, t, w, int(int(h) * STACK_MEAN_BAND_FRAC), mean_lines,
                              theme=theme)
    else:
        chart_t, chart_h = t, h

    _clone_chart(slide, tmpl["stacked"], XL_CHART_TYPE.COLUMN_STACKED_100, l, chart_t, w, chart_h, cd,
                series_colors=series_colors,
                hide_zero_labels=True, legend_n=len(choices))


# ── Slide builder ─────────────────────────────────────────────────────────────

def _build_slide(prs, layout, q, page_num, tmpl, *, title_suffix: str = "") -> None:
    slide = prs.slides.add_slide(layout)
    chart_type = q.get("chart_type", "bar_vertical")
    breakdowns = q.get("breakdowns", [])
    n_groups   = len(breakdowns)
    note_lines = [f"Question type: {_question_type_label(q)}"]

    q_label    = _shorten_label(q.get("label", ""))
    slide_title = q.get("title") or q_label
    title_text = f"{slide_title} {title_suffix}".strip() if title_suffix else slide_title

    _add_rect(slide, FOOTER_L, FOOTER_T, FOOTER_W, FOOTER_H, C_ORANGE)
    title_font = 18 if len(title_text) > 120 else 22 if len(title_text) > 80 else 30
    _add_text(slide, TITLE_L, TITLE_T, TITLE_W, TITLE_H,
              title_text, font_size=title_font, bold=True, color=C_NAVY)

    if chart_type == "donut_stacked":
        _add_section_label(slide, DT_TOTAL_L, DT_TOTAL_T, DT_TOTAL_W, "Total")
        _build_donut(slide, q, tmpl, DT_LEFT_L, DT_LEFT_T, DT_LEFT_W, DT_LEFT_H)
        if breakdowns:
            combined = " / ".join(bd["group_label"] for bd in breakdowns)
            _add_section_label(slide, DT_RIGHT_LBL_L, DT_RIGHT_LBL_T,
                               DT_RIGHT_LBL_W, combined)
            _build_stacked(slide, q, breakdowns, tmpl,
                           DT_RIGHT_L, DT_RIGHT_T, DT_RIGHT_W, DT_RIGHT_H)
    elif chart_type == "bar_horizontal":
        # MA: Total + up to MAX_GROUPS_PER_SLIDE breakdowns, 3 equal-width columns.
        _add_section_label(slide, _bh_col_l(0), BH_LBL_T, BH_COL_W, "Total")
        _build_total_bar(slide, q, tmpl, _bh_col_l(0), BH_CHART_T, BH_COL_W, BH_CHART_H)
        for gi, bd in enumerate(breakdowns[:MAX_GROUPS_PER_SLIDE]):
            col_l = _bh_col_l(gi + 1)
            _add_section_label(slide, col_l, BH_LBL_T, BH_COL_W, bd["group_label"])
            _build_breakdown_bar(slide, q, bd, tmpl,
                                 col_l, BH_CHART_T, BH_COL_W, BH_CHART_H)
        if breakdowns and len(q.get("choices", [])) >= MA_CROSSTAB_MIN_ITEMS:
            note_lines.append(MA_CROSSTAB_NOTE)

    else:
        _add_section_label(slide, BR_TOTAL_L, BR_TOTAL_T, BR_TOTAL_W, "Total")
        _build_total_col(slide, q, tmpl, BR_LEFT_L, BR_LEFT_T, BR_LEFT_W, BR_LEFT_H)

        if breakdowns:
            total_span = int(BR_RIGHT_BOTTOM - BR_RIGHT_FIRST_T)
            per_group  = total_span // max(n_groups, 1)
            chart_h    = per_group - int(SECTION_TO_CHART)
            for gi, bd in enumerate(breakdowns):
                lbl_t   = int(BR_RIGHT_FIRST_T) + gi * per_group
                chart_t = lbl_t + int(SECTION_TO_CHART)
                _add_section_label(slide, BR_RIGHT_LBL_L, Emu(lbl_t),
                                   BR_RIGHT_LBL_W, bd["group_label"])
                _build_breakdown_col(slide, q, bd, tmpl,
                                         BR_RIGHT_L, Emu(chart_t),
                                         BR_RIGHT_W, Emu(chart_h))

    q_base = q.get("total", {}).get("base", "XX")
    _add_q_label(slide, Q_LABEL_L, Q_LABEL_T, Q_LABEL_W, Q_LABEL_H,
                 q.get("question", ""), q_label, q_base)
    _set_slide_note(slide, "\n".join(note_lines))
    _add_text(slide, PAGE_L, PAGE_T, PAGE_W, PAGE_H,
              str(page_num), font_size=8, color=C_QBASE,
              align=PP_ALIGN.RIGHT, anchor="ctr")


# ── "default"-format slide builder ────────────────────────────────────────────
#
# Unlike the general format (which builds every slide from scratch — a plain
# add_shape rectangle for the footer bar, manual textboxes for title/Q-label/
# page number — see _build_slide above), "default"-format slides are built on
# the bundled appendix_templates/default_template.pptx's own "use_dz" slide
# layout, so all of the company-branded background bars/colors come for free
# from that layout/master and never need to be redrawn in code. Only the
# TWO placeholders add_slide() does NOT auto-instantiate from a layout —
# sldNum (idx 10) and footer (idx 11), both literal <p:sp> shapes in the
# layout's own spTree rather than "special" auto-generated ones — need to be
# cloned in manually per slide; see _default_find_ph / _default_new_slide.

def _default_find_ph(layout_element, idx: int):
    """Find the <p:sp> for placeholder `idx` directly in a slide layout's
    (or slide's) own spTree — used for sldNum/footer, which python-pptx's
    add_slide() does not auto-copy from the layout (only ph types with no
    special `type` attribute, e.g. this format's idx 12/13/14 "body" ones, get
    auto-instantiated)."""
    spTree = layout_element.find(qn("p:cSld")).find(qn("p:spTree"))
    for sp in spTree.findall(qn("p:sp")):
        ph = sp.find(qn("p:nvSpPr") + "/" + qn("p:nvPr") + "/" + qn("p:ph"))
        if ph is not None and ph.get("idx") == str(idx):
            return sp
    return None


def _default_set_first_slide_num(prs, start_page: int) -> None:
    """This format's page number is a real auto-updating <a:fld type="slidenum">
    field (cloned from the layout, see _default_new_slide) rather than the
    general format's plain page-number textbox — it always reflects the
    slide's actual position in the deck. To honor a non-default --start-page
    the same way the general format does, set the presentation-level "first
    slide number" (PowerPoint's Design > Slide Size > Number slides from),
    which shifts every auto field's displayed value by the same offset."""
    if start_page != 1:
        prs.part._element.set("firstSlideNum", str(start_page))


def _default_new_slide(prs, layout, sldnum_sp, ftr_sp):
    """Add a new slide from the "default"-format "use_dz" layout: idx 12/13/14
    placeholders come for free via add_slide(); idx 14 (the big unused
    "Area for charting" content placeholder — "default"-format slides place
    charts manually, like the general format) is removed; idx 10 (page number) and
    11 (footer) are cloned in from the layout's own template `<p:sp>`
    elements (captured once by the caller — see generate()).

    The clones' shape ids are reassigned to values unused on this slide —
    the layout's own template `<p:sp>` elements carry hardcoded ids (e.g.
    sldNum id="3") that collide with whatever id add_slide() happened to
    auto-assign to idx 12/13/14 on THIS slide (e.g. the title placeholder
    also getting id="3") — a real duplicate-id bug that was rendering as
    an overlapping/duplicated title textbox in PowerPoint (two shapes
    sharing one id makes rendering behavior undefined)."""
    slide = prs.slides.add_slide(layout)
    spTree = slide.shapes._spTree
    for shp in list(slide.shapes):
        if shp.is_placeholder and shp.placeholder_format.idx == 14:
            shp._element.getparent().remove(shp._element)
    next_id = max((shp.shape_id for shp in slide.shapes), default=1) + 1
    for sp in (sldnum_sp, ftr_sp):
        if sp is None:
            continue
        clone = copy.deepcopy(sp)
        clone.find(qn("p:nvSpPr")).find(qn("p:cNvPr")).set("id", str(next_id))
        next_id += 1
        spTree.append(clone)
    return slide


def _default_set_footer(slide, question: str, label: str, base) -> None:
    """Footer placeholder (idx 11): one plain gray run — "Q17_R1.  {label}
    (N=249)" — unlike the general format's 3-color _add_q_label, matching
    the bundled template's own single-run footer styling."""
    ph = slide.placeholders[11]
    tf = ph.text_frame
    tf.text = f"{question}.  {label}  (N={base})"
    run = tf.paragraphs[0].runs[0]
    run.font.color.rgb = DEFAULT_C_GRAY
    _set_run_font(run, DEFAULT_FONT)


def _build_slide_default(prs, layout, q, page_num, tmpl, sldnum_sp, ftr_sp,
                         section_label: str, *, title_suffix: str = "") -> None:
    slide = _default_new_slide(prs, layout, sldnum_sp, ftr_sp)
    theme = DEFAULT_THEME
    chart_type = q.get("chart_type", "bar_vertical")
    breakdowns = q.get("breakdowns", [])
    note_lines = [f"Question type: {_question_type_label(q)}"]

    q_label    = _shorten_label(q.get("label", ""))
    slide_title = q.get("title") or q_label
    title_text = f"{slide_title} {title_suffix}".strip() if title_suffix else slide_title

    slide.placeholders[12].text_frame.text = section_label
    slide.placeholders[13].text_frame.text = title_text

    if chart_type == "donut_stacked":
        _add_section_label(slide, DEFAULT_DT_TOTAL_L, DEFAULT_DT_TOTAL_T, DEFAULT_DT_TOTAL_W,
                           "Total", height=Emu(219456), theme=theme)
        _build_donut(slide, q, tmpl, DEFAULT_DT_LEFT_L, DEFAULT_DT_LEFT_T,
                    DEFAULT_DT_LEFT_W, DEFAULT_DT_LEFT_H, theme=theme)
        if breakdowns:
            combined = " / ".join(bd["group_label"] for bd in breakdowns)
            _add_section_label(slide, DEFAULT_DT_RIGHT_LBL_L, DEFAULT_DT_RIGHT_LBL_T,
                               DEFAULT_DT_RIGHT_LBL_W, combined,
                               height=Emu(219456), theme=theme)
            _build_stacked(slide, q, breakdowns, tmpl,
                           DEFAULT_DT_RIGHT_L, DEFAULT_DT_RIGHT_T,
                           DEFAULT_DT_RIGHT_W, DEFAULT_DT_RIGHT_H, theme=theme)
    else:
        # bar_horizontal (MA) AND the general format's bar_vertical fallback
        # (large-choice SA-Nominal) both map to "default"-format's one
        # non-donut layout — the bundled template only demonstrates a
        # horizontal-bar 3-column style, not a vertical-column one (see the
        # module comment above _default_find_ph / the "'default' format
        # geometry" constants block).
        _add_section_label(slide, _default_bh_col_l(0), DEFAULT_BH_LBL_T, DEFAULT_BH_COL_W,
                           "Total", height=Emu(219456), theme=theme)
        _build_total_bar(slide, q, tmpl, _default_bh_col_l(0), DEFAULT_BH_CHART_T,
                         DEFAULT_BH_COL_W, DEFAULT_BH_CHART_H, theme=theme)
        for gi, bd in enumerate(breakdowns[:MAX_GROUPS_PER_SLIDE]):
            col_l = _default_bh_col_l(gi + 1)
            _add_section_label(slide, col_l, DEFAULT_BH_LBL_T, DEFAULT_BH_COL_W,
                               bd["group_label"], height=Emu(219456), theme=theme)
            _build_breakdown_bar(slide, q, bd, tmpl,
                                 col_l, DEFAULT_BH_CHART_T, DEFAULT_BH_COL_W, DEFAULT_BH_CHART_H,
                                 theme=theme)
        if breakdowns and len(q.get("choices", [])) >= MA_CROSSTAB_MIN_ITEMS:
            note_lines.append(MA_CROSSTAB_NOTE)

    q_base = q.get("total", {}).get("base", "XX")
    _default_set_footer(slide, q.get("question", ""), q_label, q_base)
    _set_slide_note(slide, "\n".join(note_lines))


# ── Template loading ──────────────────────────────────────────────────────────

def _load_templates(templates_dir: str, *,
                    roles: tuple[str, ...] = ("bar", "col", "donut", "stacked"),
                    extractor_hint: str = "extract_chart_templates.py"):
    """Load the bundled chart-style chartSpace elements from XML files —
    4 roles (bar/col/donut/stacked) for the general format, 3
    (bar/donut/stacked, no "col") for "default" — see `roles`."""
    base = Path(templates_dir)
    tmpl = {}
    for role in roles:
        path = base / f"{role}.xml"
        if not path.exists():
            raise RuntimeError(
                f"Chart template '{path}' missing. Run "
                f"{extractor_hint} to regenerate.")
        tmpl[role] = parse_xml(path.read_bytes())
    return tmpl


def _general_templates_dir() -> str:
    """Return the bundled chart_templates directory inside the installed package."""
    return str(Path(__file__).parent / "chart_templates")


def _default_templates_dir() -> str:
    """Return the bundled "default"-format chart_templates directory."""
    return str(Path(__file__).parent / "chart_templates_default")


def _default_pptx_path() -> str:
    """Return the bundled (slide-stripped) "default"-format base presentation
    — see tools/extract_chart_templates_default.py's module docstring for why
    it has no slides."""
    return str(Path(__file__).parent / "appendix_templates" / "default_template.pptx")


def _default_load_base(pptx_path: str):
    """Open the bundled "default"-format base presentation and return
    (prs, use_dz_layout, sldnum_sp_template, ftr_sp_template) — the last two
    are deep-copyable <p:sp> elements for the sldNum/footer placeholders,
    captured from the layout itself since add_slide() doesn't auto-copy
    those two ph types (see _default_new_slide)."""
    prs = Presentation(pptx_path)
    layout = next((l for l in prs.slide_layouts if l.name == "use_dz"), None)
    if layout is None:
        raise RuntimeError(
            f"Slide layout 'use_dz' not found in '{pptx_path}' — the bundled "
            "\"default\"-format template may be corrupted or was replaced incorrectly.")
    layout_el = layout._element
    sldnum_sp = _default_find_ph(layout_el, 10)
    ftr_sp    = _default_find_ph(layout_el, 11)
    return prs, layout, sldnum_sp, ftr_sp


# Default customer-logo choice — the bundled template already ships with the
# Acecook logo placed top-right of the Q&Me one (see _default_apply_customer_logo).
DEFAULT_LOGO = "acecook"


def _default_apply_customer_logo(prs, logo: str | None) -> None:
    """Swap or remove the customer-logo picture that sits to the right of the
    Q&Me logo on every "default"-format slide (both are baked into
    the bundled template's slide master, so they render on every slide for
    free — see the "CustomerLogo"/"QMeLogo" shape names set on
    default_template.pptx). The Q&Me logo itself is never touched.

    `logo` (case-insensitive):
        "acecook" (default) — no-op, keep the bundled Acecook picture as-is.
        "none"              — delete the customer-logo picture AND the
                               vertical divider line between it and the Q&Me
                               logo ("LogoDivider" — otherwise a stray "|"
                               would be left dangling with nothing after it),
                               leaving only the Q&Me logo.
        anything else       — treated as a path to a replacement logo image;
                               the Acecook picture is removed and the new
                               image is inserted in the same slot, scaled
                               (preserving its own aspect ratio) to fit
                               inside the original picture's bounding box
                               and centered within it. The divider line is
                               kept — there are still two logos to separate.
    """
    logo = (logo or DEFAULT_LOGO).strip()
    if logo.lower() == DEFAULT_LOGO:
        return

    master = prs.slide_masters[0]
    pic_shape = next((s for s in master.shapes if s.name == "CustomerLogo"), None)
    if pic_shape is None:
        return  # template has no customer-logo slot to touch

    box_l, box_t = pic_shape.left, pic_shape.top
    box_w, box_h = pic_shape.width, pic_shape.height
    pic_shape._element.getparent().remove(pic_shape._element)

    if logo.lower() == "none":
        divider = next((s for s in master.shapes if s.name == "LogoDivider"), None)
        if divider is not None:
            divider._element.getparent().remove(divider._element)
        return

    logo_path = Path(logo)
    if not logo_path.exists():
        raise FileNotFoundError(f"Customer logo image not found: {logo_path}")

    image_part, rId = master.part.get_or_add_image_part(str(logo_path))
    native_w, native_h = image_part.scale(None, None)
    fit = min(int(box_w) / int(native_w), int(box_h) / int(native_h))
    final_w = int(native_w * fit)
    final_h = int(native_h * fit)
    final_x = int(box_l) + (int(box_w) - final_w) // 2
    final_y = int(box_t) + (int(box_h) - final_h) // 2
    master.shapes._spTree.add_pic(
        master.shapes._next_shape_id, "CustomerLogo", logo_path.name,
        rId, final_x, final_y, final_w, final_h,
    )


# ── Public API ────────────────────────────────────────────────────────────────

def _safe_print(msg: str) -> None:
    try:
        print(msg)
    except UnicodeEncodeError:
        print(msg.encode("ascii", errors="replace").decode("ascii"))


def _group_name(tbl: dict) -> str:
    """Display/Section name for an appendix table — its own ``sub_title``
    verbatim (no more "Appendix - " prefix-stripping: which table renders is
    now decided purely by the explicit ``is_appendix`` field, so sub_title is
    free-text with no naming convention attached to it). Falls back to
    "Table {table_index}" if sub_title is empty, so a group is never
    unlabeled."""
    st = str(tbl.get("sub_title") or "").strip()
    return st or f"Table {tbl.get('table_index', '')}"


def _auto_match_groups(tables: list[dict]) -> list[tuple[dict, str]]:
    """Every table whose ``is_appendix`` field is true, as
    ``[(table, group_name), ...]`` in chart_data.json order — see
    ``_group_name``. Does not raise; callers decide what an empty result
    means."""
    return [(tbl, _group_name(tbl)) for tbl in tables if tbl.get("is_appendix")]


def list_groups(chart_data_path: str) -> list[dict]:
    """List EVERY table in chart_data.json — table_index, sub_title, and its
    current ``is_appendix`` flag — regardless of whether it's marked yet.

    Used two ways (see CLAUDE.md Workflow D):
      1. Candidate discovery — before anything is marked, a caller (the
         Claude skill) lists every table here to ask the user which one(s)
         should become appendix tables, then sets ``"is_appendix": true`` on
         the chosen table(s) directly in datatable.json (persisted — no
         need to re-ask on later runs, same pattern as appendix_format/
         appendix_logo).
      2. Preview — on a later run, the ``is_appendix`` values already reflect
         what a plain ``generate()`` call (no ``table_idx``/``table_indices``
         override) would render, without actually generating a .pptx.

    Returns ``[{"table_index", "sub_title", "is_appendix"}, ...]`` in
    chart_data.json order (empty list only if chart_data.json has no tables
    at all).
    """
    data = json.loads(Path(chart_data_path).read_text(encoding="utf-8"))
    return [
        {
            "table_index": tbl.get("table_index"),
            "sub_title":   tbl.get("sub_title", ""),
            "is_appendix": bool(tbl.get("is_appendix")),
        }
        for tbl in data.get("tables", [])
    ]


def _select_tables(tables: list[dict], table_idx: int | None,
                   table_indices: list[int] | None = None) -> list[tuple[dict, str]]:
    """Pick which chart_data.json table(s) to render into the appendix, as
    ``[(table, group_name), ...]`` in chart_data.json order.

    Explicit ``table_idx`` always wins (CLI ``--table N`` override) and
    bypasses grouping entirely — returns exactly that one table with an empty
    group name (no PowerPoint Section).

    Otherwise, if ``table_indices`` is given (CLI ``--tables 1,3``), render
    only those ``table_index`` values — but ONLY among the tables marked
    ``"is_appendix": true`` (see ``list_groups``); this is how a caller runs
    a user-picked SUBSET of groups after asking "which table(s)?" (CLAUDE.md
    Workflow D) rather than either exactly-one (``table_idx``) or all of
    them. Raises ValueError if any requested index isn't marked. Still
    grouped into PowerPoint Sections the same way as the "all" case when
    more than one is picked.

    Otherwise auto-select EVERY table with ``"is_appendix": true`` — see
    ``_auto_match_groups``. Which tables render is decided ENTIRELY by this
    explicit field, never by ``sub_title`` naming (no more "Appendix -"/
    "General" prefix convention). Each matched table is its own group (1
    table = 1 group; datatable.json already carries one full table config
    per brand/segment, see CLAUDE.md Workflow D). Raises ValueError if
    nothing matches — appendix generation never silently renders every
    table in the file, only the ones opted in.
    """
    if table_idx is not None:
        for tbl in tables:
            if tbl.get("table_index") == table_idx:
                return [(tbl, "")]
        raise ValueError(f"No table with table_index={table_idx} in chart_data.json")

    matched = _auto_match_groups(tables)

    if table_indices is not None:
        wanted = list(dict.fromkeys(table_indices))  # de-dupe, preserve request order
        by_idx = {tbl.get("table_index"): (tbl, group) for tbl, group in matched}
        missing = [i for i in wanted if i not in by_idx]
        if missing:
            available = sorted(by_idx)
            raise ValueError(
                f"table_indices {missing} not marked \"is_appendix\": true "
                f"(available table_index values: {available}). Use "
                "list_groups() / --list-groups to see every table's current "
                "is_appendix state."
            )
        # Keep chart_data.json order among the requested subset, not request order.
        wanted_set = set(wanted)
        return [(tbl, group) for tbl, group in matched if tbl.get("table_index") in wanted_set]

    if not matched:
        available = ", ".join(repr(t.get("sub_title", "")) for t in tables) or "(none)"
        raise ValueError(
            "No table marked \"is_appendix\": true found in chart_data.json "
            f"(available sub_title values: {available}). Add "
            "\"is_appendix\": true to a table item in datatable.json (see "
            "CLAUDE.md Workflow D)."
        )
    return matched


# PowerPoint's native "Sections" feature (Slide Sorter / Outline view) is
# stored as an undocumented-by-python-pptx but standard OOXML extension:
# <p:presentation>/<p:extLst>/<p:ext uri="{521415D9-...}">/<p14:sectionLst>.
# Unknown extensions are spec-required to be ignored gracefully by any
# reader, so even if some PowerPoint build doesn't render sections from
# this, the deck still opens and reads fine — the slides themselves are
# unaffected either way.
_P14_NS = "http://schemas.microsoft.com/office/powerpoint/2010/main"
_SECTION_EXT_URI = "{521415D9-36F7-43E2-AB2F-B90AF26B5E84}"


def _add_pptx_sections(prs, sections: list[tuple[str, int, int]]) -> None:
    """Group slides into native PowerPoint Sections, so a multi-group
    appendix deck can be jumped between per-group in Slide Sorter/Outline —
    this is the ONLY grouping cue in the deck. An earlier version also
    inserted a title-only divider slide before each group, but that cost
    every group an extra page and was dropped by request — the Section
    label already carries the group name without spending a slide on it.

    `sections` is ``[(name, first_slide_idx, last_slide_idx), ...]`` — 0-based,
    inclusive, in slide-add order. Slides are only ever appended (never
    reordered), so that index order always matches `<p:sldId>` order in the
    presentation, which is what section membership is keyed on."""
    sld_ids = [sld_id.get("id") for sld_id in prs._element.sldIdLst]

    ext_lst = prs._element.find(qn("p:extLst"))
    if ext_lst is None:
        ext_lst = etree.SubElement(prs._element, qn("p:extLst"))
    ext = etree.SubElement(ext_lst, qn("p:ext"))
    ext.set("uri", _SECTION_EXT_URI)
    section_lst = etree.SubElement(ext, f"{{{_P14_NS}}}sectionLst")

    for name, first_idx, last_idx in sections:
        section = etree.SubElement(section_lst, f"{{{_P14_NS}}}section")
        section.set("name", name)
        section.set("id", "{%s}" % str(uuid.uuid4()).upper())
        sld_id_lst14 = etree.SubElement(section, f"{{{_P14_NS}}}sldIdLst")
        for idx in range(first_idx, last_idx + 1):
            etree.SubElement(sld_id_lst14, f"{{{_P14_NS}}}sldId").set("id", sld_ids[idx])


def _chunk_breakdowns_by_col_count(breakdowns: list, max_cols: int) -> list[list]:
    """Split a list of banner groups across as few slides as possible
    (never splitting one group's own columns across two slides), with each
    slide's total column count as EVEN as possible — rather than greedily
    filling each slide to the brim before starting the next (which can
    produce a lopsided e.g. 9-and-4 split), this computes the minimum
    slide count (`ceil(total / max_cols)`) and distributes whole groups
    across that many slides using LPT (longest-processing-time-first) bin
    balancing: largest groups placed first, each into whichever slide
    currently has the fewest columns — e.g. Area(4)/Gender(2)/Age(3)/HHI(4)
    = 13 cols → 2 slides balances to 7/6 instead of greedy's 9/4."""
    total = sum(len(bd.get("columns", [])) for bd in breakdowns)
    if total <= max_cols:
        return [breakdowns]

    n_slides = -(-total // max_cols)  # ceil(total / max_cols)
    order = sorted(range(len(breakdowns)),
                    key=lambda i: -len(breakdowns[i].get("columns", [])))
    bucket_idxs: list[list[int]] = [[] for _ in range(n_slides)]
    bucket_totals = [0] * n_slides
    for i in order:
        j = min(range(n_slides), key=lambda k: bucket_totals[k])
        bucket_idxs[j].append(i)
        bucket_totals[j] += len(breakdowns[i].get("columns", []))

    chunks = [
        [breakdowns[i] for i in sorted(idxs)]
        for idxs in bucket_idxs if idxs
    ]
    chunks.sort(key=lambda chunk: breakdowns.index(chunk[0]))
    return chunks


def generate(chart_data_path: str, output_path: str, *,
             templates_dir: str | None = None, table_idx: int | None = None,
             table_indices: list[int] | None = None,
             start_page: int = 1, style: str | None = None,
             default_pptx: str | None = None,
             default_section_label: str | None = None,
             logo: str | None = None) -> None:
    """Generate a PowerPoint appendix from chart_data.json.

    Renders EVERY table marked ``"is_appendix": true`` as its own group —
    see ``_select_tables`` for the selection/grouping rule — combined into a
    single deck. When more
    than one table matches, the groups are wired up as native PowerPoint
    Sections (visible/navigable in Slide Sorter and Outline view; see
    ``_add_pptx_sections``) — no extra divider slide is inserted, only the
    group's own question slides. A single matching table (the common case)
    renders exactly as before — no sections either. ``table_idx`` overrides
    all of this and renders exactly one table, ungrouped.

    Args:
        chart_data_path: Path to chart_data.json produced by the pipeline.
        output_path:     Where to write the .pptx file.
        templates_dir:   Override chart template XML directory (default:
                          bundled — chart_templates/ for "general",
                          chart_templates_default/ for "default").
        table_idx:       If set, export ONLY this table index (0-based),
                          ungrouped, instead of auto-selecting/grouping by
                          ``is_appendix``. Takes priority over ``table_indices``.
        table_indices:   If set, render only these ``table_index`` values —
                          each still its own group (a PowerPoint Section when
                          more than one) — instead of every table marked
                          ``"is_appendix": true``. Every value must already be
                          marked (see ``list_groups``); this is how a caller
                          renders a user-picked SUBSET of groups after asking
                          "which table(s)?" (CLAUDE.md Workflow D) rather than
                          exactly-one or all of them.
        start_page:      Starting page number shown in slides (default: 1),
                          continuous across every group.
        style:           "default" (company-branded, the default) or
                          "general" (surveyflow's plain style, kept as an
                          opt-in override — never auto-selected). Applies to
                          the WHOLE deck (a single .pptx can't mix styles) —
                          if not given, falls back to the FIRST selected
                          table's own ``appendix_format`` field in
                          chart_data.json (see CLAUDE.md Workflow D), then
                          "default". The legacy value "dzung_team" is
                          accepted as an alias for "default" (old
                          chart_data.json/datatable.json files).
        default_pptx:     Override the bundled "default"-format base .pptx
                          (default: bundled).
        default_section_label: Override the "default" format's language-dependent
                          section tag (default: "PHỤ LỤC"/"APPENDIX" per the
                          chart_data.json's own "lang" field — see
                          _default_section_label).
        logo:             Customer logo for the "default" format only — see
                          _default_apply_customer_logo. If not given, falls back to
                          the FIRST selected table's own ``appendix_logo`` field in
                          chart_data.json, then "acecook". Ignored when
                          style == "general" (that layout has no logos).
    """
    data = json.loads(Path(chart_data_path).read_text(encoding="utf-8"))
    selected = _select_tables(data.get("tables", []), table_idx, table_indices)
    first_tbl = selected[0][0]

    style = (style or first_tbl.get("appendix_format") or "default").strip().lower()
    if style == "dzung_team":
        style = "default"
    if style not in ("general", "default"):
        raise ValueError(f"Unknown appendix style '{style}' — expected 'general' or 'default'.")

    # A single deck can only use one visual style — the first selected
    # table's (or the CLI/API override's) style governs every group. Flag any
    # other group that explicitly asked for a different one so a mismatch
    # isn't silently swallowed.
    for tbl, group_name in selected[1:]:
        other = str(tbl.get("appendix_format") or "").strip().lower()
        other = "default" if other == "dzung_team" else other
        if other and other != style:
            _safe_print(
                f"  [warn] group '{group_name}' requests appendix_format="
                f"{other!r} but the deck-wide style is {style!r} (from the "
                "first selected table, or --format) — ignoring."
            )

    if style == "default":
        logo = logo or first_tbl.get("appendix_logo") or DEFAULT_LOGO
        _generate_default(selected, output_path, templates_dir=templates_dir,
                          start_page=start_page, default_pptx=default_pptx,
                          section_label=default_section_label,
                          lang=data.get("lang", "vi"),
                          logo=logo)
    else:
        _generate_general(selected, output_path, templates_dir=templates_dir,
                          start_page=start_page)


def _generate_general(selected: list[tuple[dict, str]], output_path: str, *,
                      templates_dir: str | None, start_page: int) -> None:
    tmpl = _load_templates(templates_dir or _general_templates_dir())

    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    blank_layout = prs.slide_layouts[6]

    def add_slide(q, page, *, title_suffix=""):
        _build_slide(prs, blank_layout, q, page, tmpl, title_suffix=title_suffix)

    multi = len(selected) > 1
    page = start_page
    sections: list[tuple[str, int, int]] = []
    for tbl, group_name in selected:
        if multi:
            _safe_print(f"\n=== Group: {group_name} ===")
        start_idx = len(prs.slides)
        page = _emit_slides(tbl, add_slide, page)
        if multi:
            sections.append((group_name, start_idx, len(prs.slides) - 1))

    if multi:
        _add_pptx_sections(prs, sections)

    prs.save(output_path)
    suffix = f", {len(selected)} groups" if multi else ""
    _safe_print(f"\nSaved -> {output_path}  ({page - start_page} slides{suffix})")


def _generate_default(selected: list[tuple[dict, str]], output_path: str, *,
                      templates_dir: str | None, start_page: int,
                      default_pptx: str | None,
                      section_label: str | None, lang: str = "vi",
                      logo: str | None = None) -> None:
    tmpl = _load_templates(templates_dir or _default_templates_dir(),
                           roles=("bar", "donut", "stacked"),
                           extractor_hint="tools/extract_chart_templates_default.py")
    prs, layout, sldnum_sp, ftr_sp = _default_load_base(default_pptx or _default_pptx_path())
    _default_set_first_slide_num(prs, start_page)
    _default_apply_customer_logo(prs, logo)
    label = section_label or _default_section_label(lang)

    def add_slide(q, page, *, title_suffix=""):
        _build_slide_default(prs, layout, q, page, tmpl, sldnum_sp, ftr_sp,
                             label, title_suffix=title_suffix)

    multi = len(selected) > 1
    page = start_page
    sections: list[tuple[str, int, int]] = []
    for tbl, group_name in selected:
        if multi:
            _safe_print(f"\n=== Group: {group_name} ===")
        start_idx = len(prs.slides)
        page = _emit_slides(tbl, add_slide, page)
        if multi:
            sections.append((group_name, start_idx, len(prs.slides) - 1))

    if multi:
        _add_pptx_sections(prs, sections)

    prs.save(output_path)
    suffix = f", {len(selected)} groups" if multi else ""
    _safe_print(f"\nSaved -> {output_path}  ({page - start_page} slides{suffix})")


def _emit_slides(tbl: dict, add_slide, start_page: int) -> int:
    """Shared question -> slide(s) dispatch loop (chunking rules for
    many-column donut_stacked / many-group bar layouts) — identical for both
    formats, since only what a single `add_slide(q, page, title_suffix=...)`
    call actually draws differs (see _generate_general / _generate_default).
    Returns the next unused page number."""
    page = start_page
    for q in tbl.get("questions", []):
        if q.get("total", {}).get("base", 0) == 0:
            continue
        _safe_print(f"  slide {page:>3}: {q.get('question', ''):<12}  {q.get('title') or q.get('label', '')}")
        chart_type = q.get("chart_type", "bar_vertical")
        breakdowns = q.get("breakdowns", [])

        if chart_type == "donut_stacked":
            total_cols = sum(len(bd.get("columns", [])) for bd in breakdowns)
            if total_cols <= STACK_MAX_COLS_PER_SLIDE:
                add_slide(q, page)
                page += 1
            else:
                # Split the stack's columns across continuation slides
                # (whole banner groups per slide, see
                # _chunk_breakdowns_by_col_count) — the Total donut is
                # re-rendered identically on every one of them.
                chunks = _chunk_breakdowns_by_col_count(breakdowns, STACK_MAX_COLS_PER_SLIDE)
                for ci, chunk in enumerate(chunks):
                    q_slide = {**q, "breakdowns": chunk}
                    add_slide(q_slide, page, title_suffix=f"({ci + 1})")
                    page += 1
        elif len(breakdowns) <= MAX_GROUPS_PER_SLIDE:
            add_slide(q, page)
            page += 1
        else:
            # Split into continuation slides, max MAX_GROUPS_PER_SLIDE groups each
            chunks = [breakdowns[i:i + MAX_GROUPS_PER_SLIDE]
                      for i in range(0, len(breakdowns), MAX_GROUPS_PER_SLIDE)]
            for ci, chunk in enumerate(chunks):
                q_slide = {**q, "breakdowns": chunk}
                add_slide(q_slide, page, title_suffix=f"({ci + 1})")
                page += 1
    return page


# ── CLI entry point ───────────────────────────────────────────────────────────

def main(argv=None) -> None:
    ap = argparse.ArgumentParser(
        description="Generate editable PowerPoint slides from chart_data.json",
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    ap.add_argument("chart_data", help="Path to chart_data.json")
    ap.add_argument("output", nargs="?", default=None,
                    help="Output .pptx path (omit with --list-groups)")
    ap.add_argument("--templates", default=None,
                    help="Chart-template dir (default: bundled package templates)")
    ap.add_argument("--list-groups", action="store_true", dest="list_groups",
                    help="Print every table in chart_data.json as JSON — "
                         "table_index, sub_title, and its current "
                         "is_appendix flag — and exit without generating a "
                         ".pptx. Use this to ask the user which table(s) "
                         "should be marked \"is_appendix\": true (CLAUDE.md "
                         "Workflow D), or to preview what a plain run would "
                         "render, or to pass table_index values to --tables.")
    ap.add_argument("--table", type=int, default=None,
                    help="Export ONLY this table index, ungrouped (default: "
                         "every table marked \"is_appendix\": true, each "
                         "rendered as its own group). Mutually exclusive "
                         "with --tables.")
    ap.add_argument("--tables", default=None,
                    help="Comma-separated table_index values to render as a "
                         "SUBSET of the tables marked \"is_appendix\": true "
                         "(each still its own group — a PowerPoint Section "
                         "when more than one), e.g. '1,3'. See "
                         "--list-groups. Mutually exclusive with --table.")
    ap.add_argument("--start-page", type=int, default=1, dest="start_page",
                    help="Starting page number (default: 1)")
    ap.add_argument("--format", choices=("general", "default"), default=None,
                    help="Appendix visual style (default: the selected table's "
                         "own \"appendix_format\" field in chart_data.json, "
                         "else 'default')")
    ap.add_argument("--default-pptx", default=None, dest="default_pptx",
                    help="Override the bundled base .pptx for the 'default' format "
                         "(--format default only)")
    ap.add_argument("--default-section-label", default=None, dest="default_section_label",
                    help="Override the 'default' format's header tag, default "
                         f"{DEFAULT_SECTION_LABEL_VI!r}/{DEFAULT_SECTION_LABEL_EN!r} "
                         "per chart_data.json's \"lang\" field (--format default only)")
    ap.add_argument("--logo", default=None,
                    help="Customer logo for the 'default' format: 'acecook' "
                         "(bundled, default), 'none' (Q&Me logo only), or a path "
                         "to a replacement logo image. Falls back to the table's "
                         "own \"appendix_logo\" field in chart_data.json, then "
                         "'acecook'. Ignored for --format general.")
    args = ap.parse_args(argv)

    if args.list_groups:
        print(json.dumps(list_groups(args.chart_data), ensure_ascii=False, indent=2))
        return

    if args.output is None:
        ap.error("the following arguments are required: output")
    if args.table is not None and args.tables:
        ap.error("--table and --tables are mutually exclusive")

    table_indices = [int(x) for x in args.tables.split(",") if x.strip()] if args.tables else None
    generate(args.chart_data, args.output,
             templates_dir=args.templates,
             table_idx=args.table, table_indices=table_indices,
             start_page=args.start_page,
             style=args.format, default_pptx=args.default_pptx,
             default_section_label=args.default_section_label,
             logo=args.logo)


if __name__ == "__main__":
    main()
