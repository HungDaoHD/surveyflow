#!/usr/bin/env python3
"""Extract the 3 "default"-format chart-style templates from a styled .pptx.

Mirrors extract_chart_templates.py's role-cloning approach (see that file for
the rationale), but for the "default" appendix format — a company-branded
slide layout ("use_dz") with only 2 demo slides / 3 distinct chart styles
(no vertical "col" role: this format's non-donut layout uses horizontal bars
for every question, see generate_pptx.py's Layout C section for why).

Writes surveyflow/steps/appendix/chart_templates_default/{bar,donut,stacked}.xml.

Roles (identified by chart type, not slide order):
    bar      bar_clustered, horizontal   → Total chart (single series, no legend;
                                            a legend is added programmatically for
                                            multi-series breakdown bars — see
                                            generate_pptx.py's _ensure_legend)
    donut    doughnut                    → Total chart (donut_stacked layout)
    stacked  col_percentStacked          → combined breakdown chart (donut_stacked layout)

Re-run this only when you want to restyle the "default"-format charts. NOTE:
the bundled appendix_templates/default_template.pptx is a STRIPPED template (0
slides — its 2 demo slides were deleted after extraction, keeping only the
"use_dz" slide master/layout that generate_pptx.py clones slides from at
runtime; see the one-off strip step in this project's session history). Keep
your own un-stripped copy of the original styled .pptx (with its 2 demo
slides + 5 charts) as the source for re-extraction — running this script
against the bundled (stripped) template will fail with "Could not find
template charts". To restyle: open your un-stripped source .pptx in
PowerPoint, adjust the 2 demo slides, save, then run:

    python tools/extract_chart_templates_default.py path/to/your_unstripped_source.pptx
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

OUT_DIR = (Path(__file__).resolve().parent.parent / "surveyflow"
           / "steps" / "appendix" / "chart_templates_default")


def _bardir(cs):
    bc = cs.find(qn("c:chart") + "/" + qn("c:plotArea") + "/" + qn("c:barChart"))
    if bc is None:
        return None
    d = bc.find(qn("c:barDir"))
    return d.get("val") if d is not None else None


def _grouping(cs):
    bc = cs.find(qn("c:chart") + "/" + qn("c:plotArea") + "/" + qn("c:barChart"))
    if bc is None:
        return None
    g = bc.find(qn("c:grouping"))
    return g.get("val") if g is not None else None


def _is_donut(cs):
    return cs.find(qn("c:chart") + "/" + qn("c:plotArea") + "/" + qn("c:doughnutChart")) is not None


def _has_legend(cs):
    return cs.find(qn("c:chart") + "/" + qn("c:legend")) is not None


def extract(pptx_path: str) -> None:
    prs = Presentation(pptx_path)
    charts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_chart:
                charts.append(shape.chart._chartSpace)

    roles = {
        # Prefer the no-legend bar chart as the "bar" template (matches the
        # general format's convention: bar.xml never ships a legend, one is
        # added programmatically for multi-series breakdown bars).
        "bar":     next((cs for cs in charts if _bardir(cs) == "bar"
                         and not _has_legend(cs)), None),
        "donut":   next((cs for cs in charts if _is_donut(cs)), None),
        "stacked": next((cs for cs in charts if _bardir(cs) == "col"
                         and _grouping(cs) == "percentStacked"), None),
    }
    missing = [k for k, v in roles.items() if v is None]
    if missing:
        sys.exit(f"Could not find template charts for roles: {missing}")

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    for role, cs in roles.items():
        path = OUT_DIR / f"{role}.xml"
        path.write_bytes(etree.tostring(cs, xml_declaration=True,
                                        encoding="UTF-8", standalone=True))
        print(f"  wrote {path}")
    print(f"\nExtracted 3 \"default\"-format chart templates -> {OUT_DIR}")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        sys.exit(
            "Usage: python tools/extract_chart_templates_default.py "
            "path/to/your_unstripped_source.pptx\n"
            "(the bundled appendix_templates/default_template.pptx has no "
            "slides left to extract from — see this file's module docstring)"
        )
    extract(sys.argv[1])
