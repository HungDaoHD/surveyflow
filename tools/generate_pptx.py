#!/usr/bin/env python3
"""Generate PowerPoint slides from chart_data.json — template-clone approach.

Instead of building charts from scratch (which never matches the template's
styling), this clones pre-extracted chart-style templates and only swaps in new
data + text. Every visual property — colors, fonts, hidden value axis,
gridlines, gap width, legend, data-label format — is inherited verbatim from the
templates, so the output looks identical to temp.pptx.

The 4 chart styles live as bundled XML in tools/chart_templates/ (extracted from
temp.pptx by extract_chart_templates.py), so the generator is self-contained —
it does NOT need temp.pptx at runtime. To restyle, edit temp.pptx in PowerPoint
and re-run extract_chart_templates.py.

Chart-style roles:
    bar      bar_clustered, horizontal   → "Total" left chart (MA)
    col      col_clustered,  vertical    → breakdown right charts / vertical Total
    donut    doughnut                    → "Total" left chart (SA≤5)
    stacked  col_percentStacked          → combined breakdown right chart

Usage:
    python tools/generate_pptx.py <chart_data.json> <output.pptx> [options]

Options:
    --templates DIR    Chart-template dir (default: tools/chart_templates)
    --table N          Export only this table index (default: all)
    --start-page N     Starting page number (default: 1)
"""

from __future__ import annotations

import argparse
import copy
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import qn
    from lxml import etree
except ImportError:
    sys.exit("Missing dependency: pip install python-pptx")


# ── Color palette (text only — chart colors come from the template) ───────────

C_NAVY  = RGBColor(0x1F, 0x4E, 0x79)  # title, section labels, Q-prefix   (#1F4E79)
C_ORANGE = RGBColor(0xED, 0x7D, 0x31)  # footer bar                        (#ED7D31)
C_QDESC = RGBColor(0x40, 0x40, 0x40)  # Q-label description                (#404040)
C_QBASE = RGBColor(0x7F, 0x7F, 0x7F)  # (N=XX) and page number            (#7F7F7F)


# ── Slide geometry — exact EMU values from temp.pptx ─────────────────────────
SW = Emu(12192000)   # 13.333"
SH = Emu(6858000)    # 7.500"

# Title
TITLE_L, TITLE_T = Emu(502920),  Emu(256032)
TITLE_W, TITLE_H = Emu(8686800), Emu(640080)

# Q-label (bottom)
Q_LABEL_L, Q_LABEL_T = Emu(502920),  Emu(5961888)
Q_LABEL_W, Q_LABEL_H = Emu(11521440), Emu(512064)

# Page number
PAGE_L, PAGE_T = Emu(11365992), Emu(6473952)
PAGE_W, PAGE_H = Emu(457200),   Emu(274320)

# Orange footer bar
FOOTER_L, FOOTER_T = Emu(0), Emu(6711696)
FOOTER_W, FOOTER_H = Emu(12188952), Emu(146304)

# Section label height ("Total", "By age")
SECTION_H = Emu(292608)    # 0.32"
SECTION_TO_CHART = Emu(228600)   # gap from section label top to chart top

# ── Layout A: donut + 100%-stacked (slide2 in template) ──────────────────────
DT_TOTAL_L, DT_TOTAL_T = Emu(548640),  Emu(1143000)
DT_TOTAL_W              = Emu(4937760)
DT_LEFT_L,  DT_LEFT_T   = Emu(548640),  Emu(1417320)
DT_LEFT_W,  DT_LEFT_H   = Emu(5120640), Emu(4023360)
DT_RIGHT_LBL_L, DT_RIGHT_LBL_T = Emu(5943600), Emu(1143000)
DT_RIGHT_LBL_W                  = Emu(5760720)
DT_RIGHT_L, DT_RIGHT_T = Emu(5897880), Emu(1417320)
DT_RIGHT_W, DT_RIGHT_H = Emu(5852160), Emu(4114800)

# ── Layout B: bar charts with N breakdown groups (slide1 in template) ─────────
BR_TOTAL_L, BR_TOTAL_T = Emu(502920),  Emu(1097280)
BR_TOTAL_W              = Emu(5394960)
BR_LEFT_L,  BR_LEFT_T   = Emu(457200),  Emu(1371600)
BR_LEFT_W,  BR_LEFT_H   = Emu(5577840), Emu(4297680)
BR_RIGHT_LBL_L = Emu(6309360)   # x of right section labels
BR_RIGHT_LBL_W = Emu(5486400)   # width of right section labels
BR_RIGHT_L     = Emu(6263640)   # x of right charts
BR_RIGHT_W     = Emu(5532120)   # width of right charts
BR_RIGHT_FIRST_T = Emu(1097280)  # y of first right section label
BR_RIGHT_BOTTOM  = Emu(5715000)  # bottom of last right chart


# ── Namespaces ────────────────────────────────────────────────────────────────
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_RID  = f"{{{_R_NS}}}id"


# ── Font / text helpers ───────────────────────────────────────────────────────

def _set_run_arial(run) -> None:
    """Set font to Arial (latin + ea + cs) on a text run."""
    run.font.name = "Arial"
    rPr = run._r.get_or_add_rPr()
    for tag, pf, cs_val in [(qn("a:ea"), "34", "-122"), (qn("a:cs"), "34", "-120")]:
        el = rPr.find(tag)
        if el is None:
            el = etree.SubElement(rPr, tag)
        el.set("typeface", "Arial")
        el.set("pitchFamily", pf)
        el.set("charset", cs_val)


def _set_txbody_margins(tf, anchor: str = "ctr") -> None:
    tf.margin_left = 0
    tf.margin_top = 0
    tf.margin_right = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    if bodyPr is not None:
        bodyPr.set("anchor", anchor)


def _add_text(slide, left, top, width, height, text, *,
              font_size: int = 12, bold: bool = False,
              color: RGBColor | None = None,
              align=PP_ALIGN.LEFT, anchor: str = "ctr") -> None:
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    _set_txbody_margins(tf, anchor)
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    _set_run_arial(run)


def _add_q_label(slide, left, top, width, height,
                 question: str, label: str, base) -> None:
    """Q-label: bold navy 'Qxx.' + gray description + gray bold (N=xx)."""
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    _set_txbody_margins(tf, anchor="t")
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.LEFT

    def _run(text, bold=False, color=C_QDESC):
        r = para.add_run()
        r.text = text
        r.font.size = Pt(10)
        r.font.bold = bold
        r.font.color.rgb = color
        _set_run_arial(r)

    _run(f"{question}.  ", bold=True, color=C_NAVY)
    _run(label)
    _run(f"  (N={base})", bold=True, color=C_QBASE)


def _add_rect(slide, left, top, width, height, fill: RGBColor) -> None:
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()


def _add_section_label(slide, left, top, width, text: str) -> None:
    _add_text(slide, left, top, width, SECTION_H, text,
              font_size=13, bold=True, color=C_NAVY,
              align=PP_ALIGN.CENTER, anchor="ctr")


# ── Chart-template cloning ────────────────────────────────────────────────────

def _force_pct_labels(chartspace) -> None:
    """Force every data-label number format to 0"%" (values injected as 0–100)."""
    for dl in chartspace.iter(qn("c:dLbls")):
        nf = dl.find(qn("c:numFmt"))
        if nf is None:
            nf = etree.Element(qn("c:numFmt"))
            dl.insert(0, nf)           # numFmt must be first child of dLbls
        nf.set("formatCode", '0"%"')
        nf.set("sourceLinked", "0")


def _remove_legend(chartspace) -> None:
    chart_el = chartspace.find(qn("c:chart"))
    if chart_el is None:
        return
    legend = chart_el.find(qn("c:legend"))
    if legend is not None:
        chart_el.remove(legend)
    pv = chart_el.find(qn("c:plotVisOnly"))
    # also flip autoTitleDeleted etc. left as-is


def _clone_chart(slide, tmpl_chartspace, pptx_type, l, t, w, h, chart_data,
                 *, drop_legend: bool = False):
    """Add a chart, replace its XML with a deep copy of the template chartSpace,
    then inject `chart_data` via replace_data() so template styling is kept."""
    gf = slide.shapes.add_chart(pptx_type, l, t, w, h, chart_data)
    cp = gf.chart_part

    # rId of the freshly-created part's embedded workbook
    new_ext = cp._element.find(qn("c:externalData"))
    new_rid = new_ext.get(_RID) if new_ext is not None else None

    cs = copy.deepcopy(tmpl_chartspace)
    t_ext = cs.find(qn("c:externalData"))
    if t_ext is not None and new_rid is not None:
        t_ext.set(_RID, new_rid)

    if drop_legend:
        _remove_legend(cs)

    # swap element, clear cached Chart/ChartWorkbook so they bind to new element
    cp._element = cs
    cp.__dict__.pop("chart", None)
    cp.__dict__.pop("chart_workbook", None)

    cp.chart.replace_data(chart_data)
    _force_pct_labels(cp._element)
    return cp.chart


def _v100(percents: dict, codes: list) -> list:
    """Fractions (0–1) → percentages (0–100) for the given choice codes."""
    return [round(percents.get(c, 0.0) * 100.0, 4) for c in codes]


# ── Per-chart builders (each clones the matching template chart) ──────────────

def _build_total_bar(slide, q, tmpl, l, t, w, h) -> None:
    """Horizontal bar, single series, sorted ascending (largest at top)."""
    choices = q["choices"]
    codes = [c["code"] for c in choices]
    labels = [c["label"] for c in choices]
    pcts = q.get("total", {}).get("percents", {})
    order = sorted(range(len(codes)), key=lambda i: pcts.get(codes[i], 0))
    labels = [labels[i] for i in order]
    codes = [codes[i] for i in order]
    cd = CategoryChartData()
    cd.categories = labels
    cd.add_series("%", _v100(pcts, codes))
    _clone_chart(slide, tmpl["bar"], XL_CHART_TYPE.BAR_CLUSTERED, l, t, w, h, cd)


def _build_total_col(slide, q, tmpl, l, t, w, h) -> None:
    """Vertical column, single series, natural order (e.g. 1–10 scale)."""
    choices = q["choices"]
    codes = [c["code"] for c in choices]
    labels = [c["label"] for c in choices]
    pcts = q.get("total", {}).get("percents", {})
    cd = CategoryChartData()
    cd.categories = labels
    cd.add_series("Total", _v100(pcts, codes))
    _clone_chart(slide, tmpl["col"], XL_CHART_TYPE.COLUMN_CLUSTERED, l, t, w, h, cd,
                 drop_legend=True)


def _build_breakdown_col(slide, q, group, tmpl, l, t, w, h, *,
                         sort_desc: bool = False) -> None:
    """Vertical clustered columns — one series per breakdown column."""
    choices = q["choices"]
    codes = [c["code"] for c in choices]
    labels = [c["label"] for c in choices]
    if sort_desc:
        total_pct = q.get("total", {}).get("percents", {})
        order = sorted(range(len(codes)),
                       key=lambda i: total_pct.get(codes[i], 0), reverse=True)
        codes = [codes[i] for i in order]
        labels = [labels[i] for i in order]
    cd = CategoryChartData()
    cd.categories = labels
    for col in group["columns"]:
        cd.add_series(col["label"], _v100(col.get("percents", {}), codes))
    _clone_chart(slide, tmpl["col"], XL_CHART_TYPE.COLUMN_CLUSTERED, l, t, w, h, cd)


def _build_donut(slide, q, tmpl, l, t, w, h) -> None:
    choices = q["choices"]
    codes = [c["code"] for c in choices]
    labels = [c["label"] for c in choices]
    pcts = q.get("total", {}).get("percents", {})
    vals = _v100(pcts, codes)
    if not any(v > 0 for v in vals):
        vals = [100.0 / len(vals)] * len(vals)
    cd = CategoryChartData()
    cd.categories = labels
    cd.add_series("Total", vals)
    _clone_chart(slide, tmpl["donut"], XL_CHART_TYPE.DOUGHNUT, l, t, w, h, cd)


def _build_stacked(slide, q, breakdowns, tmpl, l, t, w, h) -> None:
    """100%-stacked columns — series per choice, categories = breakdown columns."""
    choices = q["choices"]
    codes = [c["code"] for c in choices]
    all_cols, xlabels = [], []
    for bd in breakdowns:
        for col in bd["columns"]:
            all_cols.append(col)
            xlabels.append(col["label"])
    if not all_cols:
        return
    cd = CategoryChartData()
    cd.categories = xlabels
    for code, choice in zip(codes, choices):
        cd.add_series(choice["label"], [round(col.get("percents", {}).get(code, 0.0) * 100.0, 4)
                                        for col in all_cols])
    _clone_chart(slide, tmpl["stacked"], XL_CHART_TYPE.COLUMN_STACKED_100, l, t, w, h, cd)


# ── Slide builder ─────────────────────────────────────────────────────────────

def _build_slide(prs, layout, q, page_num, tmpl) -> None:
    slide = prs.slides.add_slide(layout)
    chart_type = q.get("chart_type", "bar_vertical")
    breakdowns = q.get("breakdowns", [])
    n_groups = len(breakdowns)

    _add_rect(slide, FOOTER_L, FOOTER_T, FOOTER_W, FOOTER_H, C_ORANGE)
    _add_text(slide, TITLE_L, TITLE_T, TITLE_W, TITLE_H,
              q.get("label", ""), font_size=30, bold=True, color=C_NAVY)

    if chart_type == "donut_stacked":
        # Layout A: donut left + one 100%-stacked chart right
        _add_section_label(slide, DT_TOTAL_L, DT_TOTAL_T, DT_TOTAL_W, "Total")
        _build_donut(slide, q, tmpl, DT_LEFT_L, DT_LEFT_T, DT_LEFT_W, DT_LEFT_H)
        if breakdowns:
            combined = " / ".join(bd["group_label"] for bd in breakdowns)
            _add_section_label(slide, DT_RIGHT_LBL_L, DT_RIGHT_LBL_T,
                               DT_RIGHT_LBL_W, combined)
            _build_stacked(slide, q, breakdowns, tmpl,
                           DT_RIGHT_L, DT_RIGHT_T, DT_RIGHT_W, DT_RIGHT_H)
    else:
        # Layout B: bar/col total left + N breakdown col charts right
        _add_section_label(slide, BR_TOTAL_L, BR_TOTAL_T, BR_TOTAL_W, "Total")
        if chart_type == "bar_horizontal":
            _build_total_bar(slide, q, tmpl, BR_LEFT_L, BR_LEFT_T, BR_LEFT_W, BR_LEFT_H)
        else:
            _build_total_col(slide, q, tmpl, BR_LEFT_L, BR_LEFT_T, BR_LEFT_W, BR_LEFT_H)

        if breakdowns:
            total_span = int(BR_RIGHT_BOTTOM - BR_RIGHT_FIRST_T)
            per_group = total_span // max(n_groups, 1)
            chart_h = per_group - int(SECTION_TO_CHART)
            for gi, bd in enumerate(breakdowns):
                lbl_t = int(BR_RIGHT_FIRST_T) + gi * per_group
                chart_t = lbl_t + int(SECTION_TO_CHART)
                _add_section_label(slide, BR_RIGHT_LBL_L, Emu(lbl_t),
                                   BR_RIGHT_LBL_W, bd["group_label"])
                _build_breakdown_col(slide, q, bd, tmpl,
                                     BR_RIGHT_L, Emu(chart_t), BR_RIGHT_W, Emu(chart_h),
                                     sort_desc=(chart_type == "bar_horizontal"))

    q_base = q.get("total", {}).get("base", "XX")
    _add_q_label(slide, Q_LABEL_L, Q_LABEL_T, Q_LABEL_W, Q_LABEL_H,
                 q.get("question", ""), q.get("label", ""), q_base)
    _add_text(slide, PAGE_L, PAGE_T, PAGE_W, PAGE_H,
              str(page_num), font_size=8, color=C_QBASE,
              align=PP_ALIGN.RIGHT, anchor="ctr")


# ── Template loading + slide deletion ─────────────────────────────────────────

def _load_templates(templates_dir: str):
    """Load the 4 bundled chart-style chartSpace elements from XML files."""
    base = Path(templates_dir)
    tmpl = {}
    for role in ("bar", "col", "donut", "stacked"):
        path = base / f"{role}.xml"
        if not path.exists():
            raise RuntimeError(
                f"Chart template '{path}' missing. Run "
                f"extract_chart_templates.py to regenerate.")
        tmpl[role] = parse_xml(path.read_bytes())
    return tmpl


# ── Main ──────────────────────────────────────────────────────────────────────

def _safe_print(msg: str) -> None:
    try:
        print(msg)
    except UnicodeEncodeError:
        print(msg.encode("ascii", errors="replace").decode("ascii"))


def generate(chart_data_path: str, output_path: str, *,
             templates_dir: str | None = None, table_idx: int | None = None,
             start_page: int = 1) -> None:
    data = json.loads(Path(chart_data_path).read_text(encoding="utf-8"))
    tmpl = _load_templates(templates_dir or _default_templates_dir())

    prs = Presentation()                 # blank deck; charts carry their own style
    prs.slide_width = SW
    prs.slide_height = SH
    blank_layout = prs.slide_layouts[6]  # "Blank" layout

    page = start_page
    for tbl in data.get("tables", []):
        if table_idx is not None and tbl.get("table_index") != table_idx:
            continue
        for q in tbl.get("questions", []):
            _safe_print(f"  slide {page:>3}: {q.get('question', ''):<12}  {q.get('label', '')}")
            _build_slide(prs, blank_layout, q, page, tmpl)
            page += 1

    prs.save(output_path)
    _safe_print(f"\nSaved -> {output_path}  ({page - start_page} slides)")


def _default_templates_dir() -> str:
    return str(Path(__file__).resolve().parent / "chart_templates")


if __name__ == "__main__":
    ap = argparse.ArgumentParser(
        description="Generate editable PowerPoint slides from chart_data.json",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__,
    )
    ap.add_argument("chart_data", help="Path to chart_data.json")
    ap.add_argument("output", help="Output .pptx path")
    ap.add_argument("--templates", default=_default_templates_dir(),
                    help="Chart-template dir (default: tools/chart_templates)")
    ap.add_argument("--table", type=int, default=None,
                    help="Export only this table index (default: all)")
    ap.add_argument("--start-page", type=int, default=1, dest="start_page",
                    help="Starting page number (default: 1)")
    args = ap.parse_args()
    generate(args.chart_data, args.output,
             templates_dir=args.templates,
             table_idx=args.table, start_page=args.start_page)
